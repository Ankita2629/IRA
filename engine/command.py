import base64
import keyword
import random
import sys
import PyPDF2
import pyttsx3
from setuptools import sic
import speech_recognition as sr
import eel
import time
from gtts import gTTS
import os
import pygame
from googletrans import Translator
import requests
import platform
import ctypes
import pyautogui
import datetime
import speedtest
import psutil
import winshell
import wikipedia
import threading
import string
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import google.generativeai as genai
import json
import re
import subprocess
import pyperclip
import feedparser
import queue
import datetime
from pathlib import Path
import json
import re
import requests
from pathlib import Path
import datetime
import base64
from io import BytesIO
from PIL import Image
import google.generativeai as genai
import keyboard
import subprocess
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import winreg
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
from ctypes import cast, POINTER
from comtypes import CLSCTX_ALL
translator = Translator()
pygame.mixer.init()
AUTO_LISTEN_MODE = True  
LISTENING_ACTIVE = False
# Global language settings
CURRENT_LANGUAGE = "en"
CURRENT_LANG_CODE = "en-in"
engine = pyttsx3.init()
engine.setProperty('rate', 200)  # Adjust 150-250 for speed
engine.setProperty('volume', 0.9)

# Language configurations
LANGUAGES = {
    "english": {"tts": "en", "stt": "en-in", "name": "English"},
    "hindi": {"tts": "hi", "stt": "hi-in", "name": "Hindi"},
    "spanish": {"tts": "es", "stt": "es-es", "name": "Spanish"},
    "french": {"tts": "fr", "stt": "fr-fr", "name": "French"},
    "german": {"tts": "de", "stt": "de-de", "name": "German"},
    "italian": {"tts": "it", "stt": "it-it", "name": "Italian"},
    "portuguese": {"tts": "pt", "stt": "pt-br", "name": "Portuguese"},
    "russian": {"tts": "ru", "stt": "ru-ru", "name": "Russian"},
    "japanese": {"tts": "ja", "stt": "ja-jp", "name": "Japanese"},
    "chinese": {"tts": "zh-cn", "stt": "zh-cn", "name": "Chinese"},
    "arabic": {"tts": "ar", "stt": "ar-sa", "name": "Arabic"},
    "korean": {"tts": "ko", "stt": "ko-kr", "name": "Korean"},
    "tamil": {"tts": "ta", "stt": "ta-in", "name": "Tamil"},
    "telugu": {"tts": "te", "stt": "te-in", "name": "Telugu"},
    "bengali": {"tts": "bn", "stt": "bn-in", "name": "Bengali"},
    "marathi": {"tts": "mr", "stt": "mr-in", "name": "Marathi"},
    "gujarati": {"tts": "gu", "stt": "gu-in", "name": "Gujarati"},
    "kannada": {"tts": "kn", "stt": "kn-in", "name": "Kannada"},
    "malayalam": {"tts": "ml", "stt": "ml-in", "name": "Malayalam"},
    "punjabi": {"tts": "pa", "stt": "pa-in", "name": "Punjabi"},
    "urdu": {"tts": "ur", "stt": "ur-pk", "name": "Urdu"},
    "thai": {"tts": "th", "stt": "th-th", "name": "Thai"},
    "vietnamese": {"tts": "vi", "stt": "vi-vn", "name": "Vietnamese"},
    "indonesian": {"tts": "id", "stt": "id-id", "name": "Indonesian"},
    "dutch": {"tts": "nl", "stt": "nl-nl", "name": "Dutch"},
    "polish": {"tts": "pl", "stt": "pl-pl", "name": "Polish"},
    "turkish": {"tts": "tr", "stt": "tr-tr", "name": "Turkish"},
}

# Configuration
TODO_FILE = "todo.txt"
WEATHER_API_KEY = "90946428f9d789855734d6b3501f9978"
reminders = []
GEMINI_API_KEY = "AIzaSyCYDb08-0XuFyK4s5EGzmmtsyieG_PjW1g"
# Global speech control flag
SPEECH_INTERRUPTED = False
genai.configure(api_key=GEMINI_API_KEY)



def speak(text, language=None, slow=False):
    """Speak text with keyboard interrupt only (Space bar)"""
    global CURRENT_LANGUAGE, SPEECH_INTERRUPTED
    
    if language and language in LANGUAGES:
        lang_code = LANGUAGES[language]["tts"]
    else:
        lang_code = CURRENT_LANGUAGE
    
    text = str(text)
    
    try:
        eel.DisplayMessage(text)
        eel.receiverText(text)
        
        # Split into sentences for better interrupt responsiveness
        sentences = text.replace('!', '.').replace('?', '.').split('.')
        sentences = [s.strip() for s in sentences if s.strip()]
        
        SPEECH_INTERRUPTED = False
        
        for idx, sentence in enumerate(sentences):
           
            if SPEECH_INTERRUPTED:
                print("Speech interrupted before sentence")
                break
            try:
                if keyboard.is_pressed('space'):
                    print("Speech interrupted by Space key")
                    SPEECH_INTERRUPTED = True
                    break
            except:
                pass  # keyboard module not available or error
            
            # Check global interrupt flag
            if SPEECH_INTERRUPTED:
                print("Speech interrupted")
                break
                
            try:
                tts = gTTS(text=sentence, lang=lang_code, slow=slow)
                filename = f"temp_audio_{idx}.mp3"
                tts.save(filename)
                
                pygame.mixer.music.load(filename)
                pygame.mixer.music.play()
                
                # Check for interrupts while playing
                while pygame.mixer.music.get_busy():
                    if SPEECH_INTERRUPTED:
                        pygame.mixer.music.stop()
                        break
                    try:
                        if keyboard.is_pressed('space'):
                            pygame.mixer.music.stop()
                            SPEECH_INTERRUPTED = True
                            break
                    except:
                        pass
                    time.sleep(0.5)
                
                pygame.mixer.music.unload()
                time.sleep(0.1)
                
                if os.path.exists(filename):
                    os.remove(filename)
                    
            except Exception as e:
                print(f"Speech error in sentence: {e}")
                continue
                
    except Exception as e:
        print(f"Speech error: {e}")
    
    # Clean up any remaining temp files
    try:
        import glob
        for temp_file in glob.glob("temp_audio_*.mp3"):
            try:
                os.remove(temp_file)
            except:
                pass
    except:
        pass


audio_queue = queue.Queue()
is_listening = False
listen_thread = None
def continuous_listener():
    """Continuously listen in background and queue audio"""
    global is_listening
    r = sr.Recognizer()
    r.energy_threshold = 4000  # Adjust based on your mic
    r.dynamic_energy_threshold = True
    r.pause_threshold = 0.8  # Shorter pause = more responsive
    
    with sr.Microphone() as source:
        print("🎤 Continuous listening started...")
        r.adjust_for_ambient_noise(source, duration=1)
        
        while is_listening:
            try:
                # Listen with shorter timeout for real-time feel
                audio = r.listen(source, timeout=1, phrase_time_limit=6)
                audio_queue.put(audio)
            except sr.WaitTimeoutError:
                continue
            except Exception as e:
                print(f"Listener error: {e}")
                time.sleep(0.5)

def start_continuous_listening():
    """Start background listening thread"""
    global is_listening, listen_thread
    
    if not is_listening:
        is_listening = True
        listen_thread = threading.Thread(target=continuous_listener, daemon=True)
        listen_thread.start()
        print("✓ Real-time listening enabled")

def stop_continuous_listening():
    """Stop background listening"""
    global is_listening
    is_listening = False
    print("✗ Real-time listening stopped")

def get_command_from_queue(timeout=2):
    """Get and process audio from queue"""
    try:
        # Wait for audio with timeout
        audio = audio_queue.get(timeout=timeout)
        
        r = sr.Recognizer()
        try:
            # Quick recognition
            query = r.recognize_google(audio, language='en-in')
            print(f"You said: {query}")
            return query.lower()
        except sr.UnknownValueError:
            return ""
        except Exception as e:
            print(f"Recognition error: {e}")
            return ""
    except queue.Empty:
        return ""

def takecommand_realtime(language=None):
    """
    IMPROVED: Takes command with real-time responsiveness
    Processes audio as soon as you speak
    """
    global CURRENT_LANG_CODE
    
    lang_code = language if language else CURRENT_LANG_CODE
    
    # Start continuous listening if not already running
    if not is_listening:
        start_continuous_listening()
    
    print(f'🎧 Listening in {lang_code}...')
    eel.DisplayMessage('Listening...')
    
    # Wait for audio from queue
    start_time = time.time()
    timeout = 10
    
    while time.time() - start_time < timeout:
        query = get_command_from_queue(timeout=1)
        if query:
            eel.DisplayMessage(query)
            time.sleep(0.5)
            return query
    
    return ""

# ALTERNATIVE: Even more responsive version with callback
def instant_recognition_callback(recognizer, audio):
    """Process audio instantly as callback"""
    try:
        query = recognizer.recognize_google(audio, language='en-in')
        print(f"⚡ Instant: {query}")
        eel.DisplayMessage(query)
        eel.senderText(query)
        
        # Process command immediately
        allCommands(query)
    except sr.UnknownValueError:
        pass
    except Exception as e:
        print(f"Instant recognition error: {e}")

def start_instant_listening():
    """
    BEST OPTION: Background listening with instant callback
    Processes speech as soon as you finish a phrase
    """
    r = sr.Recognizer()
    r.energy_threshold = 4000
    r.dynamic_energy_threshold = True
    r.pause_threshold = 0.8
    
    with sr.Microphone() as source:
        print("🎤 Instant recognition active - speak anytime!")
        r.adjust_for_ambient_noise(source, duration=1)
        
        # Start background listening
        stop_listening = r.listen_in_background(
            source, 
            instant_recognition_callback,
            phrase_time_limit=6
        )
        
        # Keep running
        try:
            while True:
                time.sleep(0.1)
        except KeyboardInterrupt:
            stop_listening(wait_for_stop=False)
            print("Stopped instant listening")    
def translate_text(text, target_lang):
    """Translate text to target language"""
    try:
        translated = translator.translate(text, dest=target_lang)
        return translated.text
    except Exception as e:
        print(f"Translation error: {e}")
        return text


# def takecommand(language=None):
#     """Take voice command in specified language"""
#     global CURRENT_LANG_CODE
    
#     lang_code = language if language else CURRENT_LANG_CODE
#     r = sr.Recognizer()

#     with sr.Microphone() as source:
#         print(f'Listening in {lang_code}....')
#         eel.DisplayMessage('Listening....')
#         r.pause_threshold = 1
#         r.adjust_for_ambient_noise(source)
        
#         try:
#             audio = r.listen(source, timeout=10, phrase_time_limit=6)
#         except sr.WaitTimeoutError:
#             print("Listening timed out")
#             return ""

#     try:
#         print('Recognizing')
#         eel.DisplayMessage('Recognizing....')
#         query = r.recognize_google(audio, language=lang_code)
#         print(f"User said: {query}")
#         eel.DisplayMessage(query)
#         time.sleep(2)
        
#     except Exception as e:
#         print(f"Recognition error: {e}")
#         return ""
    
#     return query.lower()


def takecommand():
    """
    OPTIMIZED: 3x faster voice recognition
    Reduced timeouts and faster processing
    """
    global CURRENT_LANG_CODE
    
    r = sr.Recognizer()
    
    # OPTIMIZATION: Aggressive settings for speed
    r.energy_threshold = 4000
    r.dynamic_energy_threshold = False  # Disable for speed
    r.pause_threshold = 0.8  # Faster response (was 0.8)
    r.phrase_time_limit = 5  # Max 5 seconds
    
    with sr.Microphone() as source:
        print(' Listening...')
        eel.DisplayMessage('Listening...')
        
        # OPTIMIZATION: Minimal ambient noise adjustment
        r.adjust_for_ambient_noise(source, duration=0.2)  # Was 0.3
        
        try:
            # OPTIMIZATION: Shorter timeout
            audio = r.listen(source, timeout=3, phrase_time_limit=5)
            
            print(' Processing...')
            eel.DisplayMessage('Processing...')
            
            # OPTIMIZATION: Immediate recognition
            query = r.recognize_google(audio, language=CURRENT_LANG_CODE)
            print(f"✓ {query}")
            eel.DisplayMessage(query)
            
            return query.lower()
            
        except sr.WaitTimeoutError:
            return ""
        except sr.UnknownValueError:
            return ""
        except Exception as e:
            print(f"Error: {e}")
            return ""
def change_language(language_name):
    """Change the assistant's language"""
    global CURRENT_LANGUAGE, CURRENT_LANG_CODE
    
    if language_name in LANGUAGES:
        CURRENT_LANGUAGE = LANGUAGES[language_name]["tts"]
        CURRENT_LANG_CODE = LANGUAGES[language_name]["stt"]
        
        confirmation = f"Language changed to {LANGUAGES[language_name]['name']}"
        translated = translate_text(confirmation, CURRENT_LANGUAGE)
        speak(translated)
        return True
    else:
        speak("Sorry, that language is not supported")
        return False


def list_available_languages():
    """List all available languages"""
    lang_list = ", ".join([LANGUAGES[lang]["name"] for lang in list(LANGUAGES.keys())[:10]])
    speak(f"Available languages include: {lang_list}, and more")
MUSIC_LIBRARY_PATH = Path.home() / "Music"
PLAYLISTS = {}
CURRENT_PLAYLIST = []
CURRENT_TRACK_INDEX = 0

def get_music_player_sessions():
    """
    Get all active music player sessions (Windows only)
    Returns list of active media players
    """
    try:
        active_players = []
        
        # Common music player process names
        music_apps = {
            "spotify.exe": "Spotify",
            "chrome.exe": "Chrome/YouTube",
            "msedge.exe": "Edge/YouTube", 
            "vlc.exe": "VLC",
            "wmplayer.exe": "Windows Media Player",
            "iTunes.exe": "iTunes",
            "musicbee.exe": "MusicBee",
            "aimp.exe": "AIMP",
            "foobar2000.exe": "Foobar2000",
            "winamp.exe": "Winamp"
        }
        
        # Check running processes
        for proc in psutil.process_iter(['name']):
            try:
                proc_name = proc.info['name'].lower()
                for app_exe, app_name in music_apps.items():
                    if app_exe.lower() in proc_name:
                        active_players.append(app_name)
            except:
                continue
        
        return list(set(active_players))  # Remove duplicates
        
    except Exception as e:
        print(f"Get music sessions error: {e}")
        return []


def control_media_key(key):
    """
    Send media key commands (Play/Pause, Next, Previous, etc.)
    Works with ANY media player
    """
    try:
        import keyboard
        
        media_keys = {
            "play": "play/pause media",
            "pause": "play/pause media",
            "next": "next track",
            "previous": "previous track",
            "stop": "stop media",
            "volume_up": "volume up",
            "volume_down": "volume down",
            "mute": "volume mute"
        }
        
        if key in media_keys:
            keyboard.press_and_release(media_keys[key])
            return True
        return False
        
    except Exception as e:
        print(f"Media key error: {e}")
        return False


def play_music():
    """Universal play command for any active music player"""
    try:
        active_players = get_music_player_sessions()
        
        if active_players:
            control_media_key("play")
            speak("Playing music")
            return True
        else:
            speak("No music player is currently running. Would you like me to open Spotify or play local music?")
            return False
            
    except Exception as e:
        print(f"Play music error: {e}")
        speak("Could not play music")
        return False


def pause_music():
    """Universal pause command"""
    try:
        control_media_key("pause")
        speak("Music paused")
        return True
    except Exception as e:
        print(f"Pause music error: {e}")
        speak("Could not pause music")
        return False


def next_track():
    """Skip to next track"""
    try:
        control_media_key("next")
        speak("Next track")
        return True
    except Exception as e:
        print(f"Next track error: {e}")
        speak("Could not skip to next track")
        return False


def previous_track():
    """Go to previous track"""
    try:
        control_media_key("previous")
        speak("Previous track")
        return True
    except Exception as e:
        print(f"Previous track error: {e}")
        speak("Could not go to previous track")
        return False


def stop_music():
    """Stop music playback"""
    try:
        control_media_key("stop")
        speak("Music stopped")
        return True
    except Exception as e:
        print(f"Stop music error: {e}")
        speak("Could not stop music")
        return False


def open_spotify():
    """Open Spotify application"""
    try:
        speak("Opening Spotify")
        
        if platform.system() == "Windows":
            # Try Windows Store version first
            try:
                os.startfile("spotify:")
                time.sleep(2)
                return True
            except:
                # Try desktop version
                spotify_paths = [
                    os.path.expanduser(r"~\AppData\Roaming\Spotify\Spotify.exe"),
                    r"C:\Program Files\Spotify\Spotify.exe",
                    r"C:\Program Files (x86)\Spotify\Spotify.exe"
                ]
                
                for path in spotify_paths:
                    if os.path.exists(path):
                        subprocess.Popen([path])
                        speak("Spotify opened")
                        return True
                
                speak("Spotify not found. Please install Spotify first.")
                return False
        else:
            subprocess.Popen(["spotify"])
            return True
            
    except Exception as e:
        print(f"Open Spotify error: {e}")
        speak("Could not open Spotify")
        return False


def play_on_spotify(song_name):
    """
    Play a song on Spotify using web browser
    (Requires Spotify to be logged in on browser)
    """
    try:
        speak(f"Searching for {song_name} on Spotify")
        
        # Open Spotify web player with search
        import webbrowser
        search_query = song_name.replace(" ", "+")
        url = f"https://open.spotify.com/search/{search_query}"
        webbrowser.open(url)
        
        speak(f"Opening Spotify web player. Click on {song_name} to play.")
        return True
        
    except Exception as e:
        print(f"Spotify play error: {e}")
        speak("Could not play on Spotify")
        return False


def play_on_youtube_music(song_name):
    """Play a song on YouTube Music"""
    try:
        speak(f"Playing {song_name} on YouTube Music")
        
        import webbrowser
        search_query = song_name.replace(" ", "+")
        url = f"https://music.youtube.com/search?q={search_query}"
        webbrowser.open(url)
        
        time.sleep(3)
        speak(f"YouTube Music opened. Playing {song_name}")
        return True
        
    except Exception as e:
        print(f"YouTube Music error: {e}")
        speak("Could not play on YouTube Music")
        return False


def scan_local_music_library():
    """
    Scan local music folder and build library
    """
    try:
        speak("Scanning music library")
        
        music_dir = MUSIC_LIBRARY_PATH
        
        if not music_dir.exists():
            speak("Music folder not found")
            return []
        
        # Supported audio formats
        audio_formats = ['.mp3', '.wav', '.flac', '.m4a', '.ogg', '.wma', '.aac']
        
        music_files = []
        
        for root, dirs, files in os.walk(music_dir):
            for file in files:
                if any(file.lower().endswith(fmt) for fmt in audio_formats):
                    full_path = os.path.join(root, file)
                    music_files.append({
                        'title': os.path.splitext(file)[0],
                        'path': full_path,
                        'format': os.path.splitext(file)[1]
                    })
        
        speak(f"Found {len(music_files)} songs in your music library")
        return music_files
        
    except Exception as e:
        print(f"Scan library error: {e}")
        speak("Could not scan music library")
        return []


def play_local_music(song_name=None):
    """
    Play local music file using pygame
    """
    try:
        music_library = scan_local_music_library()
        
        if not music_library:
            speak("No music files found in your Music folder")
            return False
        
        if song_name:
            # Search for specific song
            song_name_lower = song_name.lower()
            matches = [song for song in music_library if song_name_lower in song['title'].lower()]
            
            if matches:
                song_to_play = matches[0]
            else:
                speak(f"Could not find {song_name}. Playing random song instead.")
                song_to_play = random.choice(music_library)
        else:
            # Play random song
            song_to_play = random.choice(music_library)
        
        speak(f"Playing {song_to_play['title']}")
        
        # Stop any currently playing music
        pygame.mixer.music.stop()
        
        # Load and play
        pygame.mixer.music.load(song_to_play['path'])
        pygame.mixer.music.play()
        
        return True
        
    except Exception as e:
        print(f"Play local music error: {e}")
        speak("Could not play local music")
        return False


def create_playlist(playlist_name):
    """Create a new playlist"""
    try:
        global PLAYLISTS
        
        if playlist_name in PLAYLISTS:
            speak(f"Playlist {playlist_name} already exists")
            return False
        
        PLAYLISTS[playlist_name] = []
        speak(f"Playlist {playlist_name} created")
        
        # Save playlists to file
        save_playlists()
        return True
        
    except Exception as e:
        print(f"Create playlist error: {e}")
        speak("Could not create playlist")
        return False


def add_to_playlist(playlist_name, song_name):
    """Add song to playlist"""
    try:
        global PLAYLISTS
        
        if playlist_name not in PLAYLISTS:
            speak(f"Playlist {playlist_name} does not exist. Creating it now.")
            create_playlist(playlist_name)
        
        music_library = scan_local_music_library()
        song_name_lower = song_name.lower()
        matches = [song for song in music_library if song_name_lower in song['title'].lower()]
        
        if matches:
            PLAYLISTS[playlist_name].append(matches[0])
            speak(f"Added {matches[0]['title']} to {playlist_name}")
            save_playlists()
            return True
        else:
            speak(f"Could not find {song_name}")
            return False
            
    except Exception as e:
        print(f"Add to playlist error: {e}")
        speak("Could not add song to playlist")
        return False


def play_playlist(playlist_name):
    """Play a saved playlist"""
    try:
        global PLAYLISTS, CURRENT_PLAYLIST, CURRENT_TRACK_INDEX
        
        if playlist_name not in PLAYLISTS:
            speak(f"Playlist {playlist_name} not found")
            return False
        
        if not PLAYLISTS[playlist_name]:
            speak(f"Playlist {playlist_name} is empty")
            return False
        
        CURRENT_PLAYLIST = PLAYLISTS[playlist_name]
        CURRENT_TRACK_INDEX = 0
        
        speak(f"Playing playlist {playlist_name}")
        
        # Play first track
        first_song = CURRENT_PLAYLIST[0]
        pygame.mixer.music.load(first_song['path'])
        pygame.mixer.music.play()
        speak(f"Now playing {first_song['title']}")
        
        return True
        
    except Exception as e:
        print(f"Play playlist error: {e}")
        speak("Could not play playlist")
        return False


def list_playlists():
    """List all available playlists"""
    try:
        global PLAYLISTS
        
        if not PLAYLISTS:
            speak("No playlists created yet")
            return False
        
        speak(f"You have {len(PLAYLISTS)} playlists:")
        for playlist_name, songs in PLAYLISTS.items():
            speak(f"{playlist_name} with {len(songs)} songs")
            time.sleep(0.3)
        
        return True
        
    except Exception as e:
        print(f"List playlists error: {e}")
        return False


def save_playlists():
    """Save playlists to file"""
    try:
        import json
        
        playlists_file = Path.home() / "Documents" / "IRA_Playlists.json"
        
        with open(playlists_file, 'w', encoding='utf-8') as f:
            json.dump(PLAYLISTS, f, indent=4)
        
        return True
        
    except Exception as e:
        print(f"Save playlists error: {e}")
        return False


def load_playlists():
    """Load playlists from file"""
    try:
        import json
        global PLAYLISTS
        
        playlists_file = Path.home() / "Documents" / "IRA_Playlists.json"
        
        if playlists_file.exists():
            with open(playlists_file, 'r', encoding='utf-8') as f:
                PLAYLISTS = json.load(f)
            return True
        
        return False
        
    except Exception as e:
        print(f"Load playlists error: {e}")
        return False


def shuffle_mode(enable=True):
    """Enable/disable shuffle mode"""
    try:
        global CURRENT_PLAYLIST
        
        if enable:
            random.shuffle(CURRENT_PLAYLIST)
            speak("Shuffle mode enabled")
        else:
            speak("Shuffle mode disabled")
        
        return True
        
    except Exception as e:
        print(f"Shuffle error: {e}")
        return False


def repeat_mode(mode="off"):
    """
    Set repeat mode
    mode: "off", "one", "all"
    """
    try:
        if mode == "one":
            speak("Repeat one track enabled")
        elif mode == "all":
            speak("Repeat all tracks enabled")
        else:
            speak("Repeat mode disabled")
        
        return True
        
    except Exception as e:
        print(f"Repeat mode error: {e}")
        return False


def what_is_playing():
    """Get currently playing song info (works with most players)"""
    try:
        active_players = get_music_player_sessions()
        
        if not active_players:
            speak("No music is currently playing")
            return None
        
        speak(f"Music is playing on {', '.join(active_players)}")
        
        # For pygame local playback
        if pygame.mixer.music.get_busy():
            if CURRENT_PLAYLIST and CURRENT_TRACK_INDEX < len(CURRENT_PLAYLIST):
                current_song = CURRENT_PLAYLIST[CURRENT_TRACK_INDEX]
                speak(f"Currently playing: {current_song['title']}")
                return current_song
        
        return None
        
    except Exception as e:
        print(f"What is playing error: {e}")
        speak("Could not determine what's playing")
        return None


def music_volume_control(action, value=None):
    """
    Control music volume
    action: "set", "increase", "decrease"
    value: volume level (0-100) for "set" action
    """
    try:
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        
        current_volume = volume.GetMasterVolumeLevelScalar()
        
        if action == "set" and value is not None:
            new_volume = value / 100.0
            volume.SetMasterVolumeLevelScalar(new_volume, None)
            speak(f"Music volume set to {value} percent")
        
        elif action == "increase":
            new_volume = min(current_volume + 0.1, 1.0)
            volume.SetMasterVolumeLevelScalar(new_volume, None)
            speak(f"Volume increased")
        
        elif action == "decrease":
            new_volume = max(current_volume - 0.1, 0.0)
            volume.SetMasterVolumeLevelScalar(new_volume, None)
            speak(f"Volume decreased")
        
        return True
        
    except Exception as e:
        print(f"Volume control error: {e}")
        speak("Could not control volume")
        return False


def change_windows_theme(theme_mode="dark"):
    """
    Change Windows theme between Dark and Light mode
    theme_mode: "dark" or "light"
    """
    try:
        if platform.system() != "Windows":
            speak("Theme changing is currently only available on Windows")
            return False
        
        # Registry path for theme settings
        registry_path = r"Software\Microsoft\Windows\CurrentVersion\Themes\Personalize"
        
        # 0 = Dark mode, 1 = Light mode
        if theme_mode.lower() in ["dark", "black", "night"]:
            apps_value = 0
            system_value = 0
            theme_name = "Dark"
        elif theme_mode.lower() in ["light", "normal", "white", "day"]:
            apps_value = 1
            system_value = 1
            theme_name = "Light"
        else:
            speak("Invalid theme. Please say dark or light")
            return False
        
        try:
            # Open registry key
            key = winreg.OpenKey(
                winreg.HKEY_CURRENT_USER,
                registry_path,
                0,
                winreg.KEY_SET_VALUE
            )
            
            # Set Apps theme (for applications)
            winreg.SetValueEx(key, "AppsUseLightTheme", 0, winreg.REG_DWORD, apps_value)
            
            # Set System theme (for taskbar, start menu, etc.)
            winreg.SetValueEx(key, "SystemUsesLightTheme", 0, winreg.REG_DWORD, system_value)
            
            winreg.CloseKey(key)
            
            speak(f"{theme_name} theme applied successfully")
            print(f"✓ Theme changed to {theme_name} mode")
            return True
            
        except WindowsError as e:
            print(f"Registry error: {e}")
            speak("Could not change theme. Administrator rights may be required")
            return False
            
    except Exception as e:
        print(f"Theme change error: {e}")
        speak("Could not change theme")
        return False


def change_accent_color(color_name="blue"):
    """
    Change Windows accent color
    Available colors: blue, red, green, purple, orange, pink, gray
    """
    try:
        if platform.system() != "Windows":
            speak("Accent color changing is only available on Windows")
            return False
        
        # Accent color hex values (BGR format for Windows Registry)
        accent_colors = {
            "blue": 0xFF8C00FF,      # Default Windows blue
            "red": 0xFF0000E8,        # Red
            "green": 0xFF00CC6A,      # Green
            "purple": 0xFFE74856,     # Purple
            "orange": 0xFF0078D7,     # Orange
            "pink": 0xFFDA1B60,       # Pink
            "gray": 0xFF7A7574,       # Gray
            "yellow": 0xFF00FFFF,     # Yellow
            "teal": 0xFF00B7C3,       # Teal
        }
        
        color_name = color_name.lower()
        
        if color_name not in accent_colors:
            speak(f"Color {color_name} not available. Available colors are: blue, red, green, purple, orange, pink, gray, yellow, teal")
            return False
        
        color_value = accent_colors[color_name]
        
        # Registry path for accent color
        registry_path = r"Software\Microsoft\Windows\DWM"
        
        try:
            key = winreg.OpenKey(
                winreg.HKEY_CURRENT_USER,
                registry_path,
                0,
                winreg.KEY_SET_VALUE
            )
            
            # Set accent color
            winreg.SetValueEx(key, "AccentColor", 0, winreg.REG_DWORD, color_value)
            winreg.SetValueEx(key, "ColorizationColor", 0, winreg.REG_DWORD, color_value)
            
            winreg.CloseKey(key)
            
            speak(f"Accent color changed to {color_name}")
            print(f"✓ Accent color changed to {color_name}")
            
            # Restart Explorer to apply changes
            restart_explorer()
            
            return True
            
        except WindowsError as e:
            print(f"Registry error: {e}")
            speak("Could not change accent color")
            return False
            
    except Exception as e:
        print(f"Accent color error: {e}")
        speak("Could not change accent color")
        return False


def enable_transparency(enable=True):
    """
    Enable or disable transparency effects in Windows
    """
    try:
        if platform.system() != "Windows":
            speak("This feature is only available on Windows")
            return False
        
        registry_path = r"Software\Microsoft\Windows\CurrentVersion\Themes\Personalize"
        
        # 1 = Enable transparency, 0 = Disable
        value = 1 if enable else 0
        
        try:
            key = winreg.OpenKey(
                winreg.HKEY_CURRENT_USER,
                registry_path,
                0,
                winreg.KEY_SET_VALUE
            )
            
            winreg.SetValueEx(key, "EnableTransparency", 0, winreg.REG_DWORD, value)
            winreg.CloseKey(key)
            
            status = "enabled" if enable else "disabled"
            speak(f"Transparency {status}")
            return True
            
        except WindowsError as e:
            print(f"Registry error: {e}")
            return False
            
    except Exception as e:
        print(f"Transparency error: {e}")
        return False


def restart_explorer():
    """
    Restart Windows Explorer to apply theme changes
    """
    try:
        speak("Restarting Windows Explorer to apply changes")
        
        # Kill Explorer
        subprocess.run(["taskkill", "/F", "/IM", "explorer.exe"], 
                      stdout=subprocess.DEVNULL, 
                      stderr=subprocess.DEVNULL)
        
        # Wait a moment
        time.sleep(1)
        
        # Restart Explorer
        subprocess.Popen("explorer.exe")
        
        time.sleep(1)
        print("✓ Windows Explorer restarted")
        return True
        
    except Exception as e:
        print(f"Restart Explorer error: {e}")
        # Explorer usually auto-restarts even if this fails
        return False


def get_current_theme():
    """
    Get the current Windows theme mode
    """
    try:
        if platform.system() != "Windows":
            return None
        
        registry_path = r"Software\Microsoft\Windows\CurrentVersion\Themes\Personalize"
        
        try:
            key = winreg.OpenKey(
                winreg.HKEY_CURRENT_USER,
                registry_path,
                0,
                winreg.KEY_READ
            )
            
            apps_theme = winreg.QueryValueEx(key, "AppsUseLightTheme")[0]
            winreg.CloseKey(key)
            
            current_theme = "Light" if apps_theme == 1 else "Dark"
            speak(f"Current theme is {current_theme} mode")
            return current_theme
            
        except WindowsError:
            return None
            
    except Exception as e:
        print(f"Get theme error: {e}")
        return None


def toggle_theme():
    """
    Toggle between dark and light theme
    """
    try:
        current = get_current_theme()
        
        if current == "Dark":
            change_windows_theme("light")
        else:
            change_windows_theme("dark")
            
    except Exception as e:
        print(f"Toggle theme error: {e}")
        speak("Could not toggle theme")


def apply_custom_theme_preset(preset_name="blue_dark"):
    """
    Apply custom theme presets with wallpaper + theme + accent
    """
    try:
        presets = {
            "blue_dark": {
                "theme": "dark",
                "accent": "blue",
                "wallpaper_category": "space"
            },
            "nature_light": {
                "theme": "light",
                "accent": "green",
                "wallpaper_category": "nature"
            },
            "sunset": {
                "theme": "dark",
                "accent": "orange",
                "wallpaper_category": "sunset"
            },
            "ocean": {
                "theme": "light",
                "accent": "teal",
                "wallpaper_category": "ocean"
            },
            "minimal": {
                "theme": "light",
                "accent": "gray",
                "wallpaper_category": "minimal"
            },
            "cyberpunk": {
                "theme": "dark",
                "accent": "purple",
                "wallpaper_category": "city"
            }
        }
        
        if preset_name not in presets:
            speak(f"Preset {preset_name} not available")
            return False
        
        preset = presets[preset_name]
        
        speak(f"Applying {preset_name} theme preset. This may take a moment.")
        
        # Apply theme
        change_windows_theme(preset["theme"])
        time.sleep(1)
        
        # Apply accent color
        change_accent_color(preset["accent"])
        time.sleep(1)
        
        # Apply wallpaper
        change_wallpaper_random(preset["wallpaper_category"])
        
        speak(f"{preset_name} preset applied successfully")
        return True
        
    except Exception as e:
        print(f"Preset error: {e}")
        speak("Could not apply preset")
        return False


def list_available_presets():
    """
    List all available theme presets
    """
    presets = [
        "blue dark - Dark theme with blue accent and space wallpaper",
        "nature light - Light theme with green accent and nature wallpaper",
        "sunset - Dark theme with orange accent and sunset wallpaper",
        "ocean - Light theme with teal accent and ocean wallpaper",
        "minimal - Light theme with gray accent and minimal wallpaper",
        "cyberpunk - Dark theme with purple accent and city wallpaper"
    ]
    
    speak("Available theme presets:")
    for preset in presets:
        speak(preset)
def change_wallpaper_from_file(image_path):
    """
    Change desktop wallpaper from a local file
    Supports Windows, Linux, and macOS
    """
    try:
        image_path = str(Path(image_path).resolve())
        
        if not os.path.exists(image_path):
            speak("Image file not found")
            return False
        
        system_os = platform.system()
        
        if system_os == "Windows":
            # Windows wallpaper change
            SPI_SETDESKWALLPAPER = 20
            result = ctypes.windll.user32.SystemParametersInfoW(
                SPI_SETDESKWALLPAPER, 
                0, 
                image_path, 
                3  # SPIF_UPDATEINIFILE | SPIF_SENDCHANGE
            )
            
            if result:
                speak("Wallpaper changed successfully")
                return True
            else:
                speak("Failed to change wallpaper")
                return False
                
        elif system_os == "Linux":
            # Linux (GNOME)
            try:
                os.system(f'gsettings set org.gnome.desktop.background picture-uri "file://{image_path}"')
                speak("Wallpaper changed successfully")
                return True
            except:
                # Try alternative for KDE
                os.system(f'qdbus org.kde.plasmashell /PlasmaShell org.kde.PlasmaShell.evaluateScript "var allDesktops = desktops();for (i=0;i<allDesktops.length;i++){{d = allDesktops[i];d.wallpaperPlugin = \'org.kde.image\';d.currentConfigGroup = Array(\'Wallpaper\', \'org.kde.image\', \'General\');d.writeConfig(\'Image\', \'file://{image_path}\')}}"')
                speak("Wallpaper changed successfully")
                return True
                
        elif system_os == "Darwin":  # macOS
            os.system(f'osascript -e \'tell application "Finder" to set desktop picture to POSIX file "{image_path}"\'')
            speak("Wallpaper changed successfully")
            return True
        else:
            speak("Wallpaper change not supported on this OS")
            return False
            
    except Exception as e:
        print(f"Wallpaper change error: {e}")
        speak("Could not change wallpaper")
        return False


def generate_gradient_wallpaper(width=1920, height=1080, category="nature"):
    """
    Generate beautiful gradient wallpaper locally - NO DOWNLOAD NEEDED
    """
    try:
        print(f"🎨 Generating {category} wallpaper...")
        
        # Create image
        img = Image.new('RGB', (width, height))
        draw = ImageDraw.Draw(img)
        
        # Color schemes based on category
        color_schemes = {
            "nature": [(34, 139, 34), (144, 238, 144), (60, 179, 113)],  # Greens
            "ocean": [(0, 105, 148), (0, 191, 255), (135, 206, 250)],  # Blues
            "sunset": [(255, 94, 77), (255, 165, 0), (255, 215, 0)],  # Orange/Yellow
            "space": [(25, 25, 112), (72, 61, 139), (138, 43, 226)],  # Purples
            "forest": [(34, 139, 34), (0, 100, 0), (85, 107, 47)],  # Dark greens
            "mountain": [(70, 130, 180), (176, 196, 222), (245, 245, 245)],  # Blue-gray
            "minimal": [(240, 240, 240), (200, 200, 200), (180, 180, 180)],  # Grays
            "fire": [(178, 34, 34), (255, 69, 0), (255, 140, 0)],  # Reds/Orange
            "random": [tuple(random.randint(50, 255) for _ in range(3)) for _ in range(3)]
        }
        
        colors = color_schemes.get(category, color_schemes["nature"])
        
        # Create smooth gradient
        for y in range(height):
            # Calculate color transition
            ratio = y / height
            
            if ratio < 0.5:
                # Transition between first and second color
                r1, g1, b1 = colors[0]
                r2, g2, b2 = colors[1]
                t = ratio * 2
            else:
                # Transition between second and third color
                r1, g1, b1 = colors[1]
                r2, g2, b2 = colors[2]
                t = (ratio - 0.5) * 2
            
            r = int(r1 + (r2 - r1) * t)
            g = int(g1 + (g2 - g1) * t)
            b = int(b1 + (b2 - b1) * t)
            
            draw.rectangle([(0, y), (width, y+1)], fill=(r, g, b))
        
        # Add some visual interest
        add_visual_effects(img, category)
        
        return img
        
    except Exception as e:
        print(f"Generate wallpaper error: {e}")
        return None


def add_visual_effects(img, category):
    """
    Add decorative elements to wallpaper
    """
    try:
        draw = ImageDraw.Draw(img, 'RGBA')
        width, height = img.size
        
        # Add circles/shapes for visual interest
        if category in ["nature", "ocean", "space"]:
            num_circles = random.randint(5, 15)
            for _ in range(num_circles):
                x = random.randint(0, width)
                y = random.randint(0, height)
                radius = random.randint(20, 100)
                
                # Semi-transparent circles
                alpha = random.randint(10, 40)
                color = (255, 255, 255, alpha)
                
                draw.ellipse(
                    [(x-radius, y-radius), (x+radius, y+radius)],
                    fill=color
                )
        
        # Add geometric patterns for minimal style
        elif category == "minimal":
            for i in range(3):
                x1 = random.randint(0, width//2)
                y1 = random.randint(0, height//2)
                x2 = x1 + random.randint(200, 500)
                y2 = y1 + random.randint(200, 500)
                
                draw.rectangle(
                    [(x1, y1), (x2, y2)],
                    fill=(255, 255, 255, 20)
                )
        
        # Apply blur for smooth effect
        img = img.filter(ImageFilter.GaussianBlur(radius=2))
        
    except Exception as e:
        print(f"Visual effects error: {e}")


def create_pattern_wallpaper(width=1920, height=1080, category="geometric"):
    """
    Create pattern-based wallpaper
    """
    try:
        print(f"🎨 Creating {category} pattern wallpaper...")
        
        img = Image.new('RGB', (width, height), color=(240, 240, 245))
        draw = ImageDraw.Draw(img)
        
        if category == "geometric":
            # Create geometric pattern
            colors = [(100, 150, 200), (150, 180, 220), (180, 200, 240)]
            
            for x in range(0, width, 100):
                for y in range(0, height, 100):
                    color = random.choice(colors)
                    size = random.randint(30, 80)
                    
                    if random.choice([True, False]):
                        # Circle
                        draw.ellipse(
                            [(x, y), (x+size, y+size)],
                            fill=color
                        )
                    else:
                        # Square
                        draw.rectangle(
                            [(x, y), (x+size, y+size)],
                            fill=color
                        )
        
        elif category == "dots":
            # Dot pattern
            base_color = (random.randint(100, 200), random.randint(100, 200), random.randint(100, 200))
            
            for x in range(0, width, 50):
                for y in range(0, height, 50):
                    radius = random.randint(5, 15)
                    draw.ellipse(
                        [(x-radius, y-radius), (x+radius, y+radius)],
                        fill=base_color
                    )
        
        return img
        
    except Exception as e:
        print(f"Pattern wallpaper error: {e}")
        return None


def change_wallpaper_random(category="random"):
    """
    Generate and set wallpaper - NO INTERNET REQUIRED
    """
    try:
        speak(f"Creating a beautiful {category} wallpaper")
        
        # Get screen resolution
        try:
            import tkinter as tk
            root = tk.Tk()
            width = root.winfo_screenwidth()
            height = root.winfo_screenheight()
            root.destroy()
        except:
            width, height = 1920, 1080
        
        print(f"Screen resolution: {width}x{height}")
        
        # Random category if needed
        if category == "random":
            categories = ["nature", "ocean", "sunset", "space", "forest", 
                         "mountain", "minimal", "fire"]
            category = random.choice(categories)
        
        # Generate wallpaper
        if category in ["geometric", "dots"]:
            img = create_pattern_wallpaper(width, height, category)
        else:
            img = generate_gradient_wallpaper(width, height, category)
        
        if img is None:
            speak("Could not generate wallpaper")
            return False
        
        # Save wallpaper
        wallpaper_dir = Path.home() / "Pictures" / "IRA_Wallpapers"
        wallpaper_dir.mkdir(parents=True, exist_ok=True)
        
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = wallpaper_dir / f"wallpaper_{category}_{timestamp}.jpg"
        
        img.save(filename, "JPEG", quality=95)
        print(f"✓ Wallpaper saved: {filename}")
        
        # Set as wallpaper
        return change_wallpaper_from_file(str(filename))
        
    except Exception as e:
        print(f"Random wallpaper error: {e}")
        import traceback
        traceback.print_exc()
        speak("Could not create wallpaper")
        return False


def change_wallpaper_from_file(image_path):
    """
    Change desktop wallpaper from a local file
    """
    try:
        image_path = str(Path(image_path).resolve())
        
        if not os.path.exists(image_path):
            speak("Image file not found")
            return False
        
        system_os = platform.system()
        
        if system_os == "Windows":
            SPI_SETDESKWALLPAPER = 20
            result = ctypes.windll.user32.SystemParametersInfoW(
                SPI_SETDESKWALLPAPER, 
                0, 
                image_path, 
                3
            )
            
            if result:
                speak("Wallpaper changed successfully")
                return True
            else:
                speak("Failed to change wallpaper")
                return False
                
        elif system_os == "Linux":
            try:
                os.system(f'gsettings set org.gnome.desktop.background picture-uri "file://{image_path}"')
                speak("Wallpaper changed successfully")
                return True
            except:
                speak("Could not change wallpaper on Linux")
                return False
                
        elif system_os == "Darwin":
            os.system(f'osascript -e \'tell application "Finder" to set desktop picture to POSIX file "{image_path}"\'')
            speak("Wallpaper changed successfully")
            return True
        else:
            speak("Wallpaper change not supported on this OS")
            return False
            
    except Exception as e:
        print(f"Wallpaper change error: {e}")
        speak("Could not change wallpaper")
        return False


def change_wallpaper_search(search_query):
    """
    Generate wallpaper based on search query
    """
    try:
        speak(f"Creating {search_query} wallpaper")
        
        # Map search terms to categories
        category_map = {
            "nature": ["nature", "tree", "plant", "flower", "garden"],
            "ocean": ["ocean", "sea", "water", "beach", "wave"],
            "sunset": ["sunset", "sunrise", "evening", "dusk"],
            "space": ["space", "star", "galaxy", "universe", "cosmos"],
            "mountain": ["mountain", "hill", "peak", "valley"],
            "minimal": ["minimal", "simple", "clean", "modern"],
            "fire": ["fire", "flame", "warm", "hot"],
        }
        
        # Find matching category
        category = "nature"  # Default
        search_lower = search_query.lower()
        
        for cat, keywords in category_map.items():
            if any(keyword in search_lower for keyword in keywords):
                category = cat
                break
        
        return change_wallpaper_random(category)
        
    except Exception as e:
        print(f"Search wallpaper error: {e}")
        speak("Could not create wallpaper")
        return False


def get_weather(city):
    """Fetch weather information for a city"""
    try:
        url = f"http://api.openweathermap.org/data/2.5/weather?q={city}&appid={WEATHER_API_KEY}&units=metric"
        response = requests.get(url, timeout=10)
        data = response.json()

        if data.get("cod") != 200:
            speak(f"Sorry, I could not find weather information for {city}")
            return

        weather_desc = data["weather"][0]["description"]
        temp = data["main"]["temp"]
        humidity = data["main"]["humidity"]
        wind_speed = data["wind"]["speed"]

        weather_report = (f"Weather in {city}: {weather_desc}. "
                          f"Temperature: {temp} degrees Celsius. "
                          f"Humidity: {humidity} percent. "
                          f"Wind speed: {wind_speed} meters per second.")
        speak(weather_report)

    except Exception as e:
        print(f"Weather error: {e}")
        speak("Sorry, I couldn't fetch the weather right now.")


def shutdown_pc():
    """Shutdown the PC safely"""
    try:
        speak("Shutting down the computer in 10 seconds. Please save your work.")
        system_os = platform.system()

        if system_os == "Windows":
            os.system("shutdown /s /t 10")
        elif system_os in ["Linux", "Darwin"]:
            os.system("shutdown -h now")
        else:
            speak("Sorry, I cannot shutdown this operating system")
    except Exception as e:
        print(f"Shutdown error: {e}")
        speak("Sorry, I could not shutdown the PC")


def lock_pc():
    """Lock the PC"""
    try:
        system_os = platform.system()
        
        if system_os == "Windows":
            ctypes.windll.user32.LockWorkStation()
            speak("PC is now locked.")
        elif system_os == "Linux":
            os.system("gnome-screensaver-command -l")
            speak("PC is now locked.")
        elif system_os == "Darwin":
            os.system("/System/Library/CoreServices/Menu\\ Extras/User.menu/Contents/Resources/CGSession -suspend")
            speak("Mac is now locked.")
        else:
            speak("Sorry, I cannot lock this operating system.")
    except Exception as e:
        print(f"Lock PC error: {e}")
        speak("Sorry, I could not lock the PC.")


def take_screenshot():
    """Take a screenshot and save it"""
    try:
        screenshots_dir = Path.home() / "Pictures" / "Screenshots"
        screenshots_dir.mkdir(parents=True, exist_ok=True)
        
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        filename = screenshots_dir / f"screenshot_{timestamp}.png"
        
        screenshot = pyautogui.screenshot()
        screenshot.save(str(filename))
        
        speak(f"Screenshot saved successfully")
    except Exception as e:
        print(f"Screenshot error: {e}")
        speak("Sorry, I could not take the screenshot.")


def get_volume_interface():
    """Get volume control interface (Windows only)"""
    try:
        from ctypes import cast, POINTER
        from comtypes import CLSCTX_ALL
        from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
        
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        return volume
    except Exception as e:
        print(f"Volume interface error: {e}")
        return None


def increase_volume():
    """Increase system volume by 10%"""
    try:
        volume = get_volume_interface()
        if volume:
            current = volume.GetMasterVolumeLevelScalar()
            new = min(current + 0.1, 1.0)
            volume.SetMasterVolumeLevelScalar(new, None)
            speak(f"Volume increased to {int(new*100)} percent")
        else:
            speak("Volume control not available")
    except Exception as e:
        print(f"Increase volume error: {e}")
        speak("Sorry, I could not increase the volume.")


def decrease_volume():
    """Decrease system volume by 10%"""
    try:
        volume = get_volume_interface()
        if volume:
            current = volume.GetMasterVolumeLevelScalar()
            new = max(current - 0.1, 0.0)
            volume.SetMasterVolumeLevelScalar(new, None)
            speak(f"Volume decreased to {int(new*100)} percent")
        else:
            speak("Volume control not available")
    except Exception as e:
        print(f"Decrease volume error: {e}")
        speak("Sorry, I could not decrease the volume.")


def mute_volume():
    """Mute system volume"""
    try:
        volume = get_volume_interface()
        if volume:
            volume.SetMute(1, None)
            speak("Volume muted")
        else:
            speak("Volume control not available")
    except Exception as e:
        print(f"Mute volume error: {e}")
        speak("Sorry, I could not mute the volume.")


def check_internet_speed():
    """Check internet speed"""
    try:
        speak("Checking internet speed. This may take a few seconds.")
        st = speedtest.Speedtest()
        st.get_best_server()
        download_speed = st.download() / 1_000_000
        upload_speed = st.upload() / 1_000_000
        ping = st.results.ping

        speak(f"Download speed: {download_speed:.2f} megabits per second. "
              f"Upload speed: {upload_speed:.2f} megabits per second. "
              f"Ping: {int(ping)} milliseconds.")
    except Exception as e:
        print(f"Internet speed error: {e}")
        speak("Sorry, I could not check the internet speed right now.")


def tell_time():
    """Speak the current time"""
    try:
        now = datetime.datetime.now()
        current_time = now.strftime("%I:%M %p")
        speak(f"The current time is {current_time}")
    except Exception as e:
        print(f"Tell time error: {e}")
        speak("Sorry, I could not tell the time right now.")


def tell_date():
    """Speak the current date"""
    try:
        now = datetime.datetime.now()
        current_date = now.strftime("%A, %d %B %Y")
        speak(f"Today's date is {current_date}")
    except Exception as e:
        print(f"Tell date error: {e}")
        speak("Sorry, I could not tell the date right now.")


def create_note(text):
    """Create a note with the given text"""
    try:
        notes_folder = Path.home() / "Documents" / "AssistantNotes"
        notes_folder.mkdir(parents=True, exist_ok=True)

        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        filename = notes_folder / f"note_{timestamp}.txt"

        filename.write_text(text, encoding="utf-8")
        speak("Note saved successfully")
    except Exception as e:
        print(f"Create note error: {e}")
        speak("Sorry, I could not create the note.")
def create_folder(folder_name):
    """Create a folder in the Documents directory"""
    try:
        base_path = Path.home() / "Documents"
        folder_path = base_path / folder_name

        if not folder_path.exists():
            folder_path.mkdir(parents=True)
            speak(f"Folder {folder_name} created successfully in Documents.")
        else:
            speak(f"Folder {folder_name} already exists in Documents.")
    except Exception as e:
        print(f"Create folder error: {e}")
        speak("Sorry, I could not create the folder.")


def delete_file(file_path):
    """Delete a file at the given path"""
    try:
        path = Path(file_path)
        if path.exists():
            path.unlink()
            speak(f"File {path.name} has been deleted successfully.")
        else:
            speak("Sorry, the file does not exist.")
    except Exception as e:
        print(f"Delete file error: {e}")
        speak("Sorry, I could not delete the file.")


def check_battery_status():
    """Check and speak battery status"""
    try:
        battery = psutil.sensors_battery()
        if battery:
            percent = battery.percent
            plugged = battery.power_plugged
            status = "charging" if plugged else "not charging"
            speak(f"The battery is at {int(percent)} percent and is currently {status}.")
        else:
            speak("Sorry, I could not detect battery information.")
    except Exception as e:
        print(f"Battery status error: {e}")
        speak("Sorry, I could not check the battery status.")


def sleepthepc():
    """Put the PC to sleep"""
    try:
        speak("Putting the computer to sleep.")
        if platform.system() == "Windows":
            os.system("rundll32.exe powrprof.dll,SetSuspendState 0,1,0")
        else:
            speak("Sleep function not available for this operating system")
    except Exception as e:
        print(f"Sleep PC error: {e}")
        speak("Sorry, I could not put the PC to sleep.")
def restartthepc():
    """Restart the PC safely"""
    try:
        speak("Restarting the computer in 10 seconds. Please save your work.")
        system_os = platform.system()

        if system_os == "Windows":
            os.system("shutdown /r /t 10")
        elif system_os in ["Linux", "Darwin"]:
            os.system("shutdown -r now")
        else:
            speak("Sorry, I cannot restart this operating system")
    except Exception as e:
        print(f"Restart error: {e}")
        speak("Sorry, I could not restart the PC")

def empty_recycle_bin():
    """Empty the Windows Recycle Bin"""
    try:
        if platform.system() != "Windows":
            speak("This feature is only available on Windows")
            return
            
        speak("Are you sure you want to empty the Recycle Bin? Please say yes or no.")
        confirmation = takecommand()
        
        if "yes" in confirmation:
            winshell.recycle_bin().empty(confirm=False, show_progress=True, sound=True)
            speak("Recycle Bin has been emptied successfully.")
        else:
            speak("Okay, Recycle Bin not emptied.")
            
    except Exception as e:
        print(f"Empty Recycle Bin error: {e}")
        speak("Sorry, I could not empty the Recycle Bin.")


def search_wikipedia(topic):
    """Search Wikipedia for a topic"""
    try:
        speak(f"Searching Wikipedia for {topic}")
        wikipedia.set_lang("en")
        summary = wikipedia.summary(topic, sentences=3)
        speak(summary)
    except wikipedia.DisambiguationError:
        speak("There are multiple results for this topic. Please be more specific.")
    except wikipedia.PageError:
        speak("Sorry, I could not find any information on that topic.")
    except Exception as e:
        print(f"Wikipedia search error: {e}")
        speak("Sorry, I could not fetch Wikipedia information.")
def adjust_brightness(command):
    """Adjust screen brightness based on voice command."""
    current = sic.get_brightness(display=0)[0]  # Get current brightness
    if "increase brightness" in command:
        new_brightness = min(current + 10, 100)  # Increase by 10%, max 100%
        sic.set_brightness(new_brightness)
        speak(f"Brightness increased to {new_brightness}%")
    elif "decrease brightness" in command:
        new_brightness = max(current - 10, 0)  # Decrease by 10%, min 0%
        sic.set_brightness(new_brightness)
        speak(f"Brightness decreased to {new_brightness}%")
    elif "set brightness" in command:
        # Example: "set brightness to 70"
        try:
            value = int(command.split("to")[-1].strip())
            value = max(0, min(value, 100))
            sic.set_brightness(value)
            speak(f"Brightness set to {value}%")
        except:
            speak("Sorry, I couldn't understand the brightness value.")

def add_task(task):
    """Add a task to the to-do list"""
    try:
        with open(TODO_FILE, "a", encoding="utf-8") as f:
            f.write(task + "\n")
        speak(f"Task added: {task}")
    except Exception as e:
        print(f"Add task error: {e}")
        speak("Sorry, I could not add the task.")


def show_tasks():
    """Show all tasks"""
    try:
        if os.path.exists(TODO_FILE):
            with open(TODO_FILE, "r", encoding="utf-8") as f:
                tasks = f.readlines()
            if tasks:
                speak("Here are your tasks:")
                for i, task in enumerate(tasks, start=1):
                    speak(f"{i}. {task.strip()}")
            else:
                speak("Your to-do list is empty.")
        else:
            speak("Your to-do list is empty.")
    except Exception as e:
        print(f"Show tasks error: {e}")
        speak("Sorry, I could not read your tasks.")


def delete_task(task_number):
    """Delete a task by its number"""
    try:
        if os.path.exists(TODO_FILE):
            with open(TODO_FILE, "r", encoding="utf-8") as f:
                tasks = f.readlines()
            if 1 <= task_number <= len(tasks):
                removed_task = tasks.pop(task_number - 1)
                with open(TODO_FILE, "w", encoding="utf-8") as f:
                    f.writelines(tasks)
                speak(f"Removed task: {removed_task.strip()}")
            else:
                speak("Invalid task number.")
        else:
            speak("Your to-do list is empty.")
    except Exception as e:
        print(f"Delete task error: {e}")
        speak("Sorry, I could not delete the task.")


def set_reminder(task, time_str):
    """Set a reminder"""
    try:
        reminder_time = datetime.datetime.strptime(time_str, "%H:%M")
        now = datetime.datetime.now()
        reminder_datetime = now.replace(hour=reminder_time.hour, minute=reminder_time.minute, second=0)
        
        if reminder_datetime < now:
            reminder_datetime += datetime.timedelta(days=1)
            
        reminders.append((task, reminder_datetime))
        speak(f"Reminder set for {task} at {time_str}")
    except Exception as e:
        print(f"Set reminder error: {e}")
        speak("Sorry, I could not set the reminder. Please use HH:MM format like 14:30")


def reminder_checker():
    """Background thread to check reminders"""
    while True:
        now = datetime.datetime.now()
        for reminder in reminders[:]:
            task, reminder_time = reminder
            if now >= reminder_time:
                speak(f"Reminder: {task}")
                reminders.remove(reminder)
        time.sleep(30)

def check_disk_usage(drive_letter="C"):
    """Check disk usage for a given drive"""
    try:
        drive_letter = drive_letter.strip().lower()

        # 🔥 Map speech variations to proper drive letters
        drive_letter_map = {
            "see": "C",
            "sea": "C",
            "c": "C",
            "c drive": "C",
            "d": "D",
            "dee": "D",
            "d drive": "D",
            "e": "E",
            "ee": "E",
            "e drive": "E"
        }

        # Replace with mapped letter (default to uppercase input)
        drive_letter = drive_letter_map.get(drive_letter, drive_letter.upper())

        if len(drive_letter) == 1 and drive_letter in string.ascii_uppercase:
            drive_path = f"{drive_letter}:\\"

            if platform.system() != "Windows":
                speak("Disk usage check is currently only available on Windows")
                return

            partitions = [p.device.upper() for p in psutil.disk_partitions()]
            if drive_path not in partitions:
                speak(f"Drive {drive_letter} is not available.")
                return

            total, used, free = shutil.disk_usage(drive_path)
            total_gb = total // (2**30)
            used_gb = used // (2**30)
            free_gb = free // (2**30)

            speak(
                f"Disk usage for drive {drive_letter}: "
                f"Total {total_gb} gigabytes, "
                f"Used {used_gb} gigabytes, "
                f"Free {free_gb} gigabytes"
            )
        else:
            speak("Invalid drive letter. Please say a valid drive like C or D")
    except Exception as e:
        print(f"Disk usage error: {e}")
        speak("Sorry, I could not check disk usage.")



def check_system_usage():
    """Check CPU and RAM usage"""
    try:
        cpu_percent = psutil.cpu_percent(interval=1)  # CPU usage %
        memory = psutil.virtual_memory()  # RAM stats
        ram_percent = memory.percent
        total_ram_gb = round(memory.total / (1024**3), 2)
        available_ram_gb = round(memory.available / (1024**3), 2)

        result = (
            f"CPU usage is at {cpu_percent} percent. "
            f"RAM usage is at {ram_percent} percent. "
            f"Total RAM {total_ram_gb} gigabytes, "
            f"Available {available_ram_gb} gigabytes."
        )
        print(result)
        speak(result)
    except Exception as e:
        print(f"System usage error: {e}")
        speak("Sorry, I could not check system usage right now.")
def get_wifi_info():
    """Get Wi-Fi network information"""
    try:
        system_os = platform.system()
        
        if system_os == "Windows":
            # Get connected Wi-Fi name
            result = os.popen('netsh wlan show interfaces').read()
            
            if "disconnected" in result.lower() or "SSID" not in result:
                speak("You are not connected to any Wi-Fi network")
                return
            
            # Extract SSID (network name)
            ssid = None
            signal = None
            for line in result.split('\n'):
                if "SSID" in line and "BSSID" not in line:
                    ssid = line.split(":")[1].strip()
                elif "Signal" in line:
                    signal = line.split(":")[1].strip()
            
            if ssid:
                info = f"You are connected to {ssid}"
                if signal:
                    info += f" with signal strength {signal}"
                speak(info)
            else:
                speak("Could not retrieve Wi-Fi information")
                
        elif system_os == "Linux":
            # For Linux systems
            result = os.popen('iwgetid -r').read().strip()
            if result:
                speak(f"You are connected to {result}")
            else:
                speak("You are not connected to any Wi-Fi network")
                
        elif system_os == "Darwin":  # macOS
            result = os.popen('/System/Library/PrivateFrameworks/Apple80211.framework/Versions/Current/Resources/airport -I | grep " SSID"').read()
            if result:
                ssid = result.split(":")[1].strip()
                speak(f"You are connected to {ssid}")
            else:
                speak("You are not connected to any Wi-Fi network")
        else:
            speak("Wi-Fi info not available for this operating system")
            
    except Exception as e:
        print(f"Wi-Fi info error: {e}")
        speak("Sorry, I could not retrieve Wi-Fi information")


def get_network_info():
    """Get detailed network information including IP addresses"""
    try:
        # Get all network interfaces
        addrs = psutil.net_if_addrs()
        stats = psutil.net_if_stats()
        
        active_networks = []
        
        for interface_name, interface_addresses in addrs.items():
            # Check if interface is up
            if interface_name in stats and stats[interface_name].isup:
                for address in interface_addresses:
                    # IPv4 address
                    if str(address.family) == 'AddressFamily.AF_INET':
                        active_networks.append({
                            'interface': interface_name,
                            'ip': address.address,
                            'netmask': address.netmask
                        })
        
        if active_networks:
            speak("Here is your network information:")
            for network in active_networks:
                # Skip loopback
                if network['ip'] != '127.0.0.1':
                    speak(f"Interface {network['interface']}: IP address {network['ip']}")
        else:
            speak("No active network connections found")
            
    except Exception as e:
        print(f"Network info error: {e}")
        speak("Sorry, I could not retrieve network information")


def show_wifi_password(network_name=None):
    """Show saved Wi-Fi password (Windows only)"""
    try:
        if platform.system() != "Windows":
            speak("This feature is only available on Windows")
            return
        
        if not network_name:
            # Get current network
            result = os.popen('netsh wlan show interfaces').read()
            for line in result.split('\n'):
                if "SSID" in line and "BSSID" not in line:
                    network_name = line.split(":")[1].strip()
                    break
        
        if not network_name:
            speak("Could not determine the network name")
            return
        
        # Get password
        command = f'netsh wlan show profile name="{network_name}" key=clear'
        result = os.popen(command).read()
        
        password = None
        for line in result.split('\n'):
            if "Key Content" in line:
                password = line.split(":")[1].strip()
                break
        
        if password:
            speak(f"The password for {network_name} is {password}")
            print(f"Password: {password}")  # Also print for security
        else:
            speak(f"Could not retrieve password for {network_name}. You may not have saved this network.")
            
    except Exception as e:
        print(f"Wi-Fi password error: {e}")
        speak("Sorry, I could not retrieve the Wi-Fi password")


def list_available_wifi():
    """List all available Wi-Fi networks"""
    try:
        if platform.system() == "Windows":
            speak("Scanning for available Wi-Fi networks. Please wait.")
            result = os.popen('netsh wlan show networks mode=bssid').read()
            
            networks = []
            for line in result.split('\n'):
                if "SSID" in line and "BSSID" not in line:
                    ssid = line.split(":")[1].strip()
                    if ssid and ssid != "":
                        networks.append(ssid)
            
            if networks:
                speak(f"Found {len(networks)} networks:")
                for i, network in enumerate(networks[:5], 1):  # Limit to first 5
                    speak(f"{i}. {network}")
                if len(networks) > 5:
                    speak(f"And {len(networks) - 5} more networks")
            else:
                speak("No Wi-Fi networks found")
        else:
            speak("This feature is currently only available on Windows")
            
    except Exception as e:
        print(f"List Wi-Fi error: {e}")
        speak("Sorry, I could not scan for Wi-Fi networks")


def find_file_by_name(filename, language="english"):
    """
    Search for a file by name in common locations
    Returns full path if found, None otherwise
    """
    try:
        # Common search locations
        search_paths = [
            Path.home() / "Documents",
            Path.home() / "Downloads",
            Path.home() / "Desktop",
            Path.home() / "Documents" / "IRA_Code_Files",
            Path.home() / "Documents" / "IRA_Presentations",
            Path.home() / "Documents" / "AssistantNotes",
            Path.home() / "Pictures",
            Path.home() / "Music",
            Path.home() / "Videos",
            Path.home(),  # Home directory
        ]
        
        # Clean filename
        filename = filename.strip().lower()
        
        # Replace spoken words with actual symbols
        filename = filename.replace(" dot ", ".")
        filename = filename.replace(" underscore ", "_")
        filename = filename.replace(" dash ", "-")
        
        found_files = []
        
        if language == "hindi":
            speak(f"{filename} को खोजा जा रहा है", language="hindi")
        else:
            speak(f"Searching for {filename}")
        
        print(f"\nSearching for: {filename}")
        print("Searching in common locations...")
        
        for search_path in search_paths:
            if search_path.exists():
                try:
                    # Search in directory and subdirectories
                    for root, dirs, files in os.walk(search_path):
                        for file in files:
                            if filename in file.lower():
                                full_path = os.path.join(root, file)
                                found_files.append(full_path)
                                print(f"Found: {full_path}")
                        
                        # Don't go too deep (max 2 levels)
                        if root.count(os.sep) - str(search_path).count(os.sep) >= 2:
                            dirs.clear()
                            
                except PermissionError:
                    continue
        
        if found_files:
            if len(found_files) == 1:
                if language == "hindi":
                    speak("फ़ाइल मिल गई", language="hindi")
                else:
                    speak("File found")
                return found_files[0]
            else:
                # Multiple files found
                if language == "hindi":
                    speak(f"{len(found_files)} फ़ाइलें मिलीं। पहली फ़ाइल उपयोग की जाएगी", language="hindi")
                else:
                    speak(f"Found {len(found_files)} files. Using the first one")
                
                # Show all found files
                for i, file_path in enumerate(found_files[:5], 1):
                    print(f"{i}. {file_path}")
                
                return found_files[0]
        else:
            if language == "hindi":
                speak(f"{filename} नहीं मिली। कृपया पूरा path बताएं", language="hindi")
            else:
                speak(f"Could not find {filename}. Please provide the full path")
            return None
            
    except Exception as e:
        print(f"File search error: {e}")
        return None


def read_text_file_smart(filename, language="english"):
    """Read text file by searching with filename only"""
    try:
        # Try to find the file
        file_path = find_file_by_name(filename, language)
        
        if not file_path:
            return False
        
        print(f"\nReading file: {file_path}")
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read().strip()
        
        if content:
            if language == "hindi":
                speak("फ़ाइल पढ़ी जा रही है", language="hindi")
            else:
                speak("Reading file")
            
            chunks = [content[i:i+2000] for i in range(0, len(content), 2000)]
            
            for chunk in chunks:
                if language == "hindi":
                    speak(chunk, language="hindi")
                else:
                    speak(chunk)
                time.sleep(0.5)
            
            if language == "hindi":
                speak("फ़ाइल पढ़ना पूरा हुआ", language="hindi")
            else:
                speak("Finished reading")
            return True
        else:
            if language == "hindi":
                speak("फ़ाइल खाली है", language="hindi")
            else:
                speak("File is empty")
            return False
            
    except Exception as e:
        print(f"Read error: {e}")
        if language == "hindi":
            speak("फ़ाइल नहीं पढ़ी जा सकी", language="hindi")
        else:
            speak("Could not read the file")
        return False


def run_generated_code(filename, language):
    """
    FIXED: Execute the generated code safely with proper error handling
    """
    try:
        speak(f"Running your {language} code now.")
        
        if language == "Python":
            # FIXED: Use absolute path and proper Python interpreter
            result = subprocess.run(
                ['python', str(filename)],
                capture_output=True,
                text=True,
                timeout=30,
                cwd=str(Path(filename).parent)  # Set working directory
            )
            
            if result.returncode == 0:
                speak("Code executed successfully!")
                if result.stdout:
                    print(f"\n=== OUTPUT ===\n{result.stdout}")
                    # Speak first 100 characters of output
                    output_preview = result.stdout[:100]
                    speak(f"Output: {output_preview}")
            else:
                speak("There was an error running the code.")
                print(f"\n=== ERROR ===\n{result.stderr}")
                # Speak the error
                error_msg = result.stderr.split('\n')[-2] if result.stderr else "Unknown error"
                speak(f"Error: {error_msg[:100]}")
                
        elif language == "JavaScript":
            try:
                result = subprocess.run(
                    ['node', str(filename)],
                    capture_output=True,
                    text=True,
                    timeout=30,
                    cwd=str(Path(filename).parent)
                )
                if result.returncode == 0:
                    speak("JavaScript code executed successfully!")
                    if result.stdout:
                        print(f"\n=== OUTPUT ===\n{result.stdout}")
                        speak(f"Output: {result.stdout[:100]}")
                else:
                    speak("Error running JavaScript code.")
                    print(f"\n=== ERROR ===\n{result.stderr}")
            except FileNotFoundError:
                speak("Node.js is not installed. Please install it to run JavaScript code.")
                print("Install Node.js from: https://nodejs.org/")
                
        else:
            speak(f"{language} code saved. Please compile and run it manually from VS Code.")
            
    except subprocess.TimeoutExpired:
        speak("Code execution timed out. It may be waiting for input or taking too long.")
    except Exception as e:
        print(f"Run error: {e}")
        import traceback
        traceback.print_exc()
        speak("Could not run the code. Check console for details.")



STABILITY_API_KEY = None 
GEMINI_API_KEY = "AIzaSyCYDb08-0XuFyK4s5EGzmmtsyieG_PjW1g"

genai.configure(api_key=GEMINI_API_KEY)


def enhance_prompt_with_ai(user_prompt):
    """
    Use Gemini to enhance and optimize the image generation prompt
    This dramatically improves image quality and accuracy
    """
    try:
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        
        enhancement_request = f"""You are an expert at writing prompts for AI image generation.
        
User's request: "{user_prompt}"

Transform this into a detailed, optimized prompt for image generation that will produce accurate, high-quality results. Include:
- Specific visual details
- Style and artistic direction
- Lighting and atmosphere
- Quality modifiers (like "highly detailed", "professional", "8k")

Return ONLY the enhanced prompt text, nothing else. Keep it under 200 words."""

        response = model.generate_content(enhancement_request)
        enhanced = response.text.strip()
        
        print(f"Original prompt: {user_prompt}")
        print(f"Enhanced prompt: {enhanced}")
        
        return enhanced
        
    except Exception as e:
        print(f"Prompt enhancement error: {e}")
        # Fallback: add basic quality modifiers
        return f"{user_prompt}, highly detailed, professional quality, sharp focus, 8k resolution"


def generate_image_pollinations_enhanced(prompt, width=1024, height=1024):
    """
    Enhanced Pollinations AI with AI-powered prompt optimization
    FREE - No API key required
    """
    try:
        # Enhance prompt using Gemini
        enhanced_prompt = enhance_prompt_with_ai(prompt)
        speak(f"Generating high-quality image. Please wait...")
        
        # Use FLUX model for better quality
        url = f"https://image.pollinations.ai/prompt/{requests.utils.quote(enhanced_prompt)}"
        params = {
            "width": width,
            "height": height,
            "nologo": "true",
            "enhance": "true",
            "model": "flux"  # Best free model
        }
        
        response = requests.get(url, params=params, timeout=90)
        
        if response.status_code == 200:
            images_dir = Path.home() / "Pictures" / "IRA_Generated_Images"
            images_dir.mkdir(parents=True, exist_ok=True)
            
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            safe_prompt = "".join(c for c in prompt[:30] if c.isalnum() or c in (' ', '-', '_')).strip()
            filename = images_dir / f"{safe_prompt}_{timestamp}.png"
            
            with open(filename, 'wb') as f:
                f.write(response.content)
            
            speak("Image generated successfully and saved to Pictures folder")
            print(f"Image saved: {filename}")
            
            try:
                import os
                os.startfile(str(filename))
            except:
                pass
                
            return True
        else:
            speak("Generation failed. Please try again.")
            return False
            
    except Exception as e:
        print(f"Image generation error: {e}")
        speak("Sorry, I could not generate the image right now.")
        return False


def generate_image_stability_ai(prompt, api_key=None):
    """
    Stability AI - BEST QUALITY (Requires API key)
    Get free credits: https://platform.stability.ai/
    """
    try:
        if not api_key:
            speak("Stability AI key not configured. Using free service instead.")
            return generate_image_pollinations_enhanced(prompt)
        
        enhanced_prompt = enhance_prompt_with_ai(prompt)
        speak("Generating professional-grade image. This may take a moment...")
        
        url = "https://api.stability.ai/v1/generation/stable-diffusion-xl-1024-v1-0/text-to-image"
        
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }
        
        data = {
            "text_prompts": [
                {
                    "text": enhanced_prompt,
                    "weight": 1
                },
                {
                    "text": "blurry, low quality, distorted, ugly, bad anatomy",
                    "weight": -1  # Negative prompt for better results
                }
            ],
            "cfg_scale": 8,  # Higher = more prompt adherence
            "height": 1024,
            "width": 1024,
            "samples": 1,
            "steps": 50,  # More steps = better quality
            "style_preset": "photographic"  # Options: photographic, digital-art, 3d-model, anime, etc.
        }
        
        response = requests.post(url, headers=headers, json=data, timeout=120)
        
        if response.status_code == 200:
            result = response.json()
            images_dir = Path.home() / "Pictures" / "IRA_Generated_Images"
            images_dir.mkdir(parents=True, exist_ok=True)
            
            for i, image in enumerate(result["artifacts"]):
                timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
                safe_prompt = "".join(c for c in prompt[:30] if c.isalnum() or c in (' ', '-', '_')).strip()
                filename = images_dir / f"{safe_prompt}_{timestamp}_HD.png"
                
                img_data = base64.b64decode(image["base64"])
                with open(filename, 'wb') as f:
                    f.write(img_data)
                
                speak("Professional-grade image generated successfully")
                print(f"Image saved: {filename}")
                
                try:
                    import os
                    os.startfile(str(filename))
                except:
                    pass
                    
            return True
        else:
            error_msg = response.json().get('message', 'Unknown error')
            print(f"API Error: {error_msg}")
            speak("API error. Switching to free service.")
            return generate_image_pollinations_enhanced(prompt)
            
    except Exception as e:
        print(f"Stability AI error: {e}")
        speak("Switching to alternative service.")
        return generate_image_pollinations_enhanced(prompt)


def generate_image_huggingface(prompt, hf_token=None):
    """
    Hugging Face - Good quality, FREE
    Optional token for faster processing: https://huggingface.co/settings/tokens
    """
    try:
        enhanced_prompt = enhance_prompt_with_ai(prompt)
        speak("Creating your image. Please wait...")
        
        # FLUX.1-schnell is fast and free
        API_URL = "https://api-inference.huggingface.co/models/black-forest-labs/FLUX.1-schnell"
        
        headers = {"Content-Type": "application/json"}
        if hf_token:
            headers["Authorization"] = f"Bearer {hf_token}"
        
        payload = {"inputs": enhanced_prompt}
        
        response = requests.post(API_URL, headers=headers, json=payload, timeout=120)
        
        if response.status_code == 200:
            images_dir = Path.home() / "Pictures" / "IRA_Generated_Images"
            images_dir.mkdir(parents=True, exist_ok=True)
            
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            safe_prompt = "".join(c for c in prompt[:30] if c.isalnum() or c in (' ', '-', '_')).strip()
            filename = images_dir / f"{safe_prompt}_{timestamp}.png"
            
            with open(filename, 'wb') as f:
                f.write(response.content)
            
            speak("Image created successfully and saved")
            print(f"Image saved: {filename}")
            
            try:
                import os
                os.startfile(str(filename))
            except:
                pass
                
            return True
        else:
            speak("Service warming up. Please try again in a moment.")
            return False
            
    except Exception as e:
        print(f"Hugging Face error: {e}")
        speak("Could not generate image.")
        return False


def generate_image_replicate(prompt):
    """
    Replicate API - Free tier available
    No API key needed 
    """
    try:
        enhanced_prompt = enhance_prompt_with_ai(prompt)
        speak("Creating your image with advanced AI. Please wait...")
        
        # Using Pollinations with better parameters as most reliable free option
        url = "https://image.pollinations.ai/prompt/" + requests.utils.quote(enhanced_prompt)
        
        params = {
            "width": 1024,
            "height": 1024,
            "seed": -1,  
            "nologo": "true",
            "enhance": "true",
            "model": "flux-pro"  
        }
        
        response = requests.get(url, params=params, timeout=90)
        
        if response.status_code == 200:
            images_dir = Path.home() / "Pictures" / "IRA_Generated_Images"
            images_dir.mkdir(parents=True, exist_ok=True)
            
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            safe_prompt = "".join(c for c in prompt[:30] if c.isalnum() or c in (' ', '-', '_')).strip()
            filename = images_dir / f"{safe_prompt}_{timestamp}.png"
            
            with open(filename, 'wb') as f:
                f.write(response.content)
            
            speak("Image created successfully")
            print(f"Image saved: {filename}")
            
            try:
                import os
                os.startfile(str(filename))
            except:
                pass
                
            return True
        return False
            
    except Exception as e:
        print(f"Replicate error: {e}")
        return False


def generate_image_multiple_attempts(prompt, service="auto"):
    """
    Smart image generation with reliable fallback options
    Removed Hugging Face due to reliability issues
    """
    try:
        enhanced_prompt = enhance_prompt_with_ai(prompt)
        
        # Try services in order of quality and reliability
        if service == "auto":
            # Try Stability AI if key available (BEST QUALITY)
            if STABILITY_API_KEY:
                print("Attempting Stability AI (best quality)...")
                result = generate_image_stability_ai(prompt, STABILITY_API_KEY)
                if result:
                    return True
                print("Stability AI failed, trying alternatives...")
            
            # Try Pollinations Enhanced (MOST RELIABLE FREE)
            print("Using Pollinations AI with enhanced prompts...")
            result = generate_image_pollinations_enhanced(prompt)
            if result:
                return True
            
            # Final fallback to Replicate
            print("Trying alternative service...")
            return generate_image_replicate(prompt)
        
        elif service == "stability":
            return generate_image_stability_ai(prompt, STABILITY_API_KEY)
        elif service == "replicate":
            return generate_image_replicate(prompt)
        else:
            return generate_image_pollinations_enhanced(prompt)
            
    except Exception as e:
        print(f"All services failed: {e}")
        speak("Unable to generate image at this time. Please try again later.")
        return False


def generate_image(prompt, service="auto", **kwargs):
    """
    Main image generation function with AI-enhanced prompts
    
    Args:
        prompt: User's description of the image
        service: 'auto' (tries all), 'stability', 'huggingface', or 'pollinations'
        **kwargs: api_key, width, height, style, etc.
    """
    return generate_image_multiple_attempts(prompt, service)


# PowerPoint Generation Functions - Fixed Version

# Configure Gemini API Key
GEMINI_API_KEY = "AIzaSyCYDb08-0XuFyK4s5EGzmmtsyieG_PjW1g"
genai.configure(api_key=GEMINI_API_KEY)


def generate_ppt_content_with_ai(topic, num_slides=7):
    """
    FIXED: Generate clean content without star symbols
    """
    try:
        speak(f"Generating content for presentation on {topic}")

        model = genai.GenerativeModel('models/gemini-2.5-flash')

        prompt = f"""
        Create a PowerPoint presentation outline on "{topic}" with exactly {num_slides} slides.
        
        IMPORTANT RULES:
        1. NO star symbols (*) anywhere
        2. NO special characters or symbols
        3. Use plain text only
        4. Each bullet point should be 10-15 words maximum
        5. No markdown formatting
        
        Return ONLY valid JSON in this exact format:

        {{
            "title": "Main Presentation Title",
            "subtitle": "Engaging subtitle",
            "slides": [
                {{
                    "title": "Introduction",
                    "content": [
                        "First clear point without symbols",
                        "Second informative point",
                        "Third relevant point"
                    ],
                    "image_keyword": "business",
                    "bg_color": "light"
                }},
                {{
                    "title": "Key Points",
                    "content": [
                        "Important detail one",
                        "Important detail two",
                        "Important detail three"
                    ],
                    "image_keyword": "technology",
                    "bg_color": "gradient"
                }}
            ]
        }}
        
        Requirements:
        - First slide: Title slide
        - Middle slides: 3-4 SHORT bullet points each (max 15 words per point)
        - Last slide: Conclusion or Thank You
        - Each slide needs "image_keyword" (one simple word like: business, people, technology, success, education, team, innovation, data, growth, future)
        - bg_color options: light, dark, gradient
        - NO stars, NO special symbols, PLAIN TEXT ONLY
        """

        response = model.generate_content(prompt)
        content_text = response.text.strip()

        # Clean markdown and special characters
        content_text = content_text.replace('```json', '').replace('```', '').strip()
        content_text = content_text.replace('*', '').replace('#', '')  # Remove stars and hashes

        try:
            content_data = json.loads(content_text)
            
            # FIXED: Clean all text content of stars and symbols
            for slide in content_data.get('slides', []):
                slide['title'] = slide['title'].replace('*', '').replace('#', '').strip()
                slide['content'] = [point.replace('*', '').replace('#', '').strip() 
                                   for point in slide['content']]
            
            print(f"✓ Generated clean content with {len(content_data.get('slides', []))} slides")
            return content_data

        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {e}")
            # Try to extract JSON
            json_match = re.search(r'\{[\s\S]*\}', content_text)
            if json_match:
                content_data = json.loads(json_match.group())
                # Clean stars from extracted JSON too
                for slide in content_data.get('slides', []):
                    slide['title'] = slide['title'].replace('*', '').strip()
                    slide['content'] = [point.replace('*', '').strip() for point in slide['content']]
                return content_data
            raise

    except Exception as e:
        print(f"AI content generation error: {e}")
        import traceback
        traceback.print_exc()
        speak("Using template content for presentation")

        # FIXED: Template without stars
        return {
            "title": topic.title(),
            "subtitle": "A Comprehensive Overview",
            "slides": [
                {"title": topic.title(), 
                 "content": ["Comprehensive overview", "Key insights", "Practical applications"], 
                 "image_keyword": "business", "bg_color": "gradient"},
                {"title": "Introduction", 
                 "content": ["Background and context", "Scope of discussion", "Importance and relevance"], 
                 "image_keyword": "people", "bg_color": "light"},
                {"title": "Key Concepts", 
                 "content": ["Core principles explained", "Fundamental theories", "Essential frameworks"], 
                 "image_keyword": "idea", "bg_color": "light"},
                {"title": "Analysis", 
                 "content": ["Data driven findings", "Expert perspectives", "Key takeaways"], 
                 "image_keyword": "data", "bg_color": "dark"},
                {"title": "Applications", 
                 "content": ["Real world examples", "Practical use cases", "Implementation strategies"], 
                 "image_keyword": "technology", "bg_color": "light"},
                {"title": "Conclusion", 
                 "content": ["Summary of main points", "Important takeaways", "Future recommendations"], 
                 "image_keyword": "success", "bg_color": "gradient"},
                {"title": "Thank You", 
                 "content": ["Questions and Discussion"], 
                 "image_keyword": "team", "bg_color": "gradient"}
            ][:num_slides]
        }


def create_enhanced_ppt(topic, num_slides=7, theme="modern"):
    """
    FIXED: Create PowerPoint with working images and no text overflow
    """
    try:
        print(f"\n{'='*60}")
        print(f"🎨 CREATING ENHANCED PRESENTATION")
        print(f"📋 Topic: {topic}")
        print(f"📊 Slides: {num_slides}")
        print(f"🎭 Theme: {theme}")
        print(f"{'='*60}\n")

        content_data = generate_ppt_content_with_ai(topic, num_slides)
        if not content_data or 'slides' not in content_data:
            speak("Failed to generate content properly")
            return False

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        themes = {
            "professional": {
                "title_color": RGBColor(0, 51, 102),
                "text_color": RGBColor(51, 51, 51),
                "accent": RGBColor(0, 112, 192),
                "gradient": "blue"
            },
            "modern": {
                "title_color": RGBColor(75, 0, 130),
                "text_color": RGBColor(60, 60, 60),
                "accent": RGBColor(138, 43, 226),
                "gradient": "purple"
            },
            "vibrant": {
                "title_color": RGBColor(200, 50, 40),
                "text_color": RGBColor(51, 51, 51),
                "accent": RGBColor(255, 100, 50),
                "gradient": "orange"
            },
            "dark": {
                "title_color": RGBColor(255, 255, 255),
                "text_color": RGBColor(220, 220, 220),
                "accent": RGBColor(100, 200, 255),
                "gradient": "dark"
            }
        }

        colors = themes.get(theme, themes["modern"])
        slides_data = content_data.get('slides', [])
        main_title = content_data.get('title', topic.title()).replace('*', '')
        subtitle = content_data.get('subtitle', f"Created by IRA • {datetime.datetime.now().strftime('%B %d, %Y')}").replace('*', '')

        # TITLE SLIDE
        blank_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(blank_layout)
        add_gradient_background(title_slide, colors["gradient"])

        # Title
        title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = main_title
        title_frame.word_wrap = True
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        subtitle_frame.paragraphs[0].font.size = Pt(20)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # CONTENT SLIDES with IMAGES
        for idx, slide_data in enumerate(slides_data[1:], 1):
            print(f"\n📄 Creating slide {idx + 1}/{len(slides_data)}...")
            slide = prs.slides.add_slide(blank_layout)
            slide_title = slide_data.get('title', f'Slide {idx + 1}').replace('*', '')
            content_points = [point.replace('*', '').strip() for point in slide_data.get('content', [])]
            image_keyword = slide_data.get('image_keyword', 'business')
            bg_style = slide_data.get('bg_color', 'light')

            # Background
            if bg_style == "gradient":
                add_gradient_background(slide, colors["gradient"])
                text_color = RGBColor(255, 255, 255)
                title_color = RGBColor(255, 255, 255)
            elif bg_style == "dark":
                add_solid_background(slide, (35, 35, 45))
                text_color = RGBColor(240, 240, 240)
                title_color = RGBColor(255, 255, 255)
            else:
                add_solid_background(slide, (250, 250, 255))
                text_color = colors["text_color"]
                title_color = colors["title_color"]

            # FIXED: Add image with proper positioning
            try:
                print(f"  📷 Downloading image for: {image_keyword}")
                image_path = download_image_multi_source(image_keyword, width=800, height=600)
                
                if image_path and image_path.exists():
                    # Add image on right side
                    left = Inches(5.5)
                    top = Inches(2)
                    width = Inches(4)
                    height = Inches(4.5)
                    
                    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
                    print(f"  ✅ Image added successfully")
                    
                    # Delete temporary image
                    try:
                        image_path.unlink()
                    except:
                        pass
                else:
                    print(f"  ⚠️ Could not download image for {image_keyword}")
            except Exception as img_error:
                print(f"  ❌ Image error: {img_error}")

            # Accent bar
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.8), Inches(0.15), Inches(0.6))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = colors["accent"]
            accent_bar.line.fill.background()

            # FIXED: Title with proper width to avoid overlap
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = slide_title
            title_frame.word_wrap = True
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = title_color

            # FIXED: Content with proper width and wrapping
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(4.5), Inches(5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            for point in content_points:
                # Remove any remaining stars
                clean_point = point.replace('*', '').strip()
                if not clean_point:
                    continue
                    
                p = text_frame.add_paragraph()
                p.text = f"• {clean_point}"
                p.font.size = Pt(16)  # Reduced from 20 to fit better
                p.font.color.rgb = text_color
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.2
                p.level = 0

        # Save presentation
        ppt_folder = Path.home() / "Documents" / "IRA_Presentations"
        ppt_folder.mkdir(parents=True, exist_ok=True)
        filename = ppt_folder / f"{topic}_{datetime.datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.pptx"
        prs.save(str(filename))

        speak(f"Presentation on {topic} created successfully with {len(prs.slides)} slides and images")
        print(f"✅ Saved: {filename}")

        try:
            os.startfile(str(filename))
        except Exception:
            pass

        return True

    except Exception as e:
        print(f"Create PPT error: {e}")
        import traceback
        traceback.print_exc()
        speak("Sorry, I encountered an error creating the presentation.")
        return False

def download_image_multi_source(keyword, width=1200, height=800):
    """Download images from multiple free sources with fallbacks"""
    keyword = keyword.split()[0].lower()

    sources = [
        {"url": f"https://picsum.photos/{width}/{height}", "name": "Picsum"},
        {"url": f"https://loremflickr.com/{width}/{height}/{keyword}", "name": "LoremFlickr"},
        {"url": f"https://via.placeholder.com/{width}x{height}/4A90E2/ffffff?text={keyword.capitalize()}", "name": "Placeholder"}
    ]

    for source in sources:
        try:
            headers = {'User-Agent': 'Mozilla/5.0'}
            response = requests.get(source["url"], timeout=10, allow_redirects=True, headers=headers)

            if response.status_code == 200 and len(response.content) > 5000:
                img = Image.open(BytesIO(response.content))
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')

                temp_path = Path(f"temp_slide_{keyword}_{random.randint(1000, 9999)}.jpg")
                img.save(temp_path, "JPEG", quality=85)
                print(f"✓ Image downloaded from {source['name']}")
                return temp_path

        except Exception as e:
            print(f"✗ {source['name']} failed: {str(e)[:50]}")
            continue

    print(f"✗ All sources failed for: {keyword}")
    return None


def add_gradient_background(slide, style="blue"):
    """Add gradient background to slide"""
    try:
        background = slide.background
        fill = background.fill
        fill.gradient()

        gradients = {
            "blue": [(25, 60, 120), (60, 120, 200)],
            "purple": [(75, 0, 130), (138, 43, 226)],
            "green": [(20, 100, 50), (50, 200, 100)],
            "orange": [(200, 80, 40), (255, 140, 60)],
            "teal": [(0, 128, 128), (64, 224, 208)],
            "dark": [(30, 30, 40), (60, 60, 80)],
            "red": [(150, 30, 30), (220, 60, 60)]
        }

        colors = gradients.get(style, gradients["blue"])
        fill.gradient_angle = 45.0
        fill.gradient_stops[0].color.rgb = RGBColor(*colors[0])
        fill.gradient_stops[1].color.rgb = RGBColor(*colors[1])

    except Exception as e:
        print(f"Gradient error: {e}")


def add_solid_background(slide, color_rgb):
    """Add solid color background"""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color_rgb)
    except Exception:
        pass


def create_ppt_from_topic(topic, num_slides=7, design_theme="modern"):
    """Wrapper for enhanced PPT creation"""
    return create_enhanced_ppt(topic, num_slides, design_theme)

# ============ CODE GENERATION CONFIGURATION ============
# Choose your preferred AI for code generation
USE_OPENAI = False  # Set to True if you have OpenAI API key
OPENAI_API_KEY = None  # Add your OpenAI key here if using

def generate_code_with_ai(code_description, language_hint=None):
    """
    Generate high-quality code using AI
    Returns: (code_string, detected_language, file_extension)
    """
    try:
        if USE_OPENAI and OPENAI_API_KEY:
            return generate_code_openai(code_description, language_hint)
        else:
            return generate_code_gemini(code_description, language_hint)
    except Exception as e:
        print(f"Code generation error: {e}")
        return None, None, None
def extract_number_from_speech(text):
    """
    FIXED: Extract numbers from spoken text more reliably
    """
    # Direct number mapping for spoken numbers
    number_words = {
        'zero': 0, 'one': 1, 'two': 2, 'three': 3, 'four': 4,
        'five': 5, 'six': 6, 'seven': 7, 'eight': 8, 'nine': 9,
        'ten': 10, 'eleven': 11, 'twelve': 12, 'thirteen': 13,
        'fourteen': 14, 'fifteen': 15, 'sixteen': 16, 'seventeen': 17,
        'eighteen': 18, 'nineteen': 19, 'twenty': 20, 'thirty': 30,
        'forty': 40, 'fifty': 50, 'sixty': 60, 'seventy': 70,
        'eighty': 80, 'ninety': 90, 'hundred': 100
    }
    
    text = text.lower().strip()
    
    # Try to extract direct digits first
    digits = ''.join(filter(str.isdigit, text))
    if digits:
        return int(digits)
    
    # Try word-to-number conversion
    for word, num in number_words.items():
        if word in text:
            return num
    
    # Try splitting and combining (e.g., "twenty five" -> 25)
    words = text.split()
    total = 0
    for word in words:
        if word in number_words:
            total += number_words[word]
    
    if total > 0:
        return total
    
    return None


def generate_code_gemini(code_description, language_hint=None):
    """
    Generate code using Google Gemini (FREE)
    """
    try:
        # UPDATED: Use the latest stable Gemini model
        model = genai.GenerativeModel('models/gemini-2.5-flash')  # Fast and stable
        
        lang_instruction = f"Write this in {language_hint}" if language_hint else "Choose the most appropriate programming language"
        
        prompt = f"""You are an expert programmer. Generate production-ready, professional code.

TASK: {code_description}

REQUIREMENTS:
- {lang_instruction}
- Include ALL necessary imports and dependencies
- Add comprehensive error handling (try-catch blocks)
- Write clear, descriptive comments
- Follow best practices and coding standards
- Make code immediately runnable
- Add input validation where needed
- Include a docstring/header comment explaining the code

OUTPUT FORMAT:
- Return ONLY executable code
- No markdown formatting
- No explanations outside code comments
- Start directly with imports or code

Generate the complete, working code now:"""

        response = model.generate_content(prompt)
        generated_code = response.text.strip()
        
        # Clean markdown if present
        generated_code = generated_code.replace('```python', '').replace('```javascript', '')
        generated_code = generated_code.replace('```java', '').replace('```cpp', '')
        generated_code = generated_code.replace('```c++', '').replace('```', '').strip()
        
        # Detect language and extension
        language, extension = detect_code_language(generated_code)
        
        print(f"Generated {language} code ({len(generated_code)} characters)")
        return generated_code, language, extension
        
    except Exception as e:
        print(f"Gemini generation error: {e}")
        import traceback
        traceback.print_exc()
        return None, None, None


def generate_code_openai(code_description, language_hint=None):
    """
    Generate code using OpenAI (PAID - Better Quality)
    """
    try:
        import openai
        openai.api_key = OPENAI_API_KEY
        
        lang_instruction = f"in {language_hint}" if language_hint else ""
        
        prompt = f"""Generate production-ready code {lang_instruction} for: {code_description}

Requirements:
- Include all imports and dependencies
- Add error handling
- Write clear comments
- Follow best practices
- Make it immediately executable

Return only the code, no explanations."""

        response = openai.ChatCompletion.create(
            model="gpt-4",  # or "gpt-3.5-turbo" for cheaper option
            messages=[
                {"role": "system", "content": "You are an expert programmer who writes clean, production-ready code."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3  # Lower = more focused, deterministic code
        )
        
        generated_code = response.choices[0].message.content.strip()
        language, extension = detect_code_language(generated_code)
        
        return generated_code, language, extension
        
    except Exception as e:
        print(f"OpenAI generation error: {e}")
        return None, None, None


def detect_code_language(code):
    """
    Detect programming language from code content
    Returns: (language_name, file_extension)
    """
    code_lower = code.lower()
    
    # Python detection
    if 'import ' in code or 'def ' in code or 'print(' in code or 'from ' in code:
        return 'Python', '.py'
    
    # JavaScript/Node.js detection
    elif 'const ' in code or 'let ' in code or 'function' in code or 'console.log' in code or 'require(' in code:
        return 'JavaScript', '.js'
    
    # Java detection
    elif 'public class' in code or 'public static void main' in code or 'System.out.println' in code:
        return 'Java', '.java'
    
    # C++ detection
    elif '#include' in code or 'std::' in code or 'cout' in code:
        return 'C++', '.cpp'
    
    # C# detection
    elif 'using System' in code or 'namespace' in code:
        return 'C#', '.cs'
    
    # HTML detection
    elif '<!DOCTYPE' in code or '<html' in code or '<body' in code:
        return 'HTML', '.html'
    
    # Default to Python
    else:
        return 'Python', '.py'


def open_vscode_and_write_code(code_description):
    """
    Main function: Generate code, save to file, open in VS Code, ask user to run
    """
    try:
        speak(f"Generating code for {code_description}. Please wait.")
        
        # Generate code using AI
        generated_code, language, extension = generate_code_with_ai(code_description)
        
        if not generated_code:
            speak("Sorry, I could not generate the code. Please try again.")
            return False
        
        # Create save directory
        code_folder = Path.home() / "Documents" / "IRA_Generated_Code"
        code_folder.mkdir(parents=True, exist_ok=True)
        
        # Create filename
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        safe_desc = "".join(c for c in code_description[:40] if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_desc = safe_desc.replace(' ', '_')
        filename = code_folder / f"{safe_desc}_{timestamp}{extension}"
        
        # Write code to file
        with open(filename, 'w', encoding='utf-8') as f:
            f.write(generated_code)
        
        speak(f"{language} code generated successfully. Opening in VS Code.")
        print(f"\nCode saved to: {filename}")
        print(f"Language: {language}")
        print(f"Lines of code: {len(generated_code.split(chr(10)))}")
        
        # Open in VS Code
        vscode_opened = open_in_vscode(filename)
        
        if not vscode_opened:
            speak("Could not open VS Code, but code is saved in Documents folder")
        
        # Wait for user to review
        time.sleep(2)
        
        # Ask user if they want to run it
        speak("Would you like me to run this code? Say yes or no.")
        confirmation = takecommand()
        
        if confirmation and "yes" in confirmation.lower():
            run_generated_code(filename, language)
        else:
            speak("Okay. The code is ready in VS Code whenever you need it.")
        
        return True
        
    except Exception as e:
        print(f"VS Code integration error: {e}")
        import traceback
        traceback.print_exc()
        speak("Sorry, I encountered an error while generating the code.")
        return False


def open_in_vscode(filepath):
    """
    Open a file in Visual Studio Code
    """
    try:
        # Try command line 'code' command
        subprocess.Popen(['code', str(filepath)], 
                        stdout=subprocess.DEVNULL, 
                        stderr=subprocess.DEVNULL)
        return True
    except FileNotFoundError:
        # Try Windows installation paths
        if platform.system() == "Windows":
            vscode_paths = [
                r"C:\Program Files\Microsoft VS Code\Code.exe",
                r"C:\Program Files (x86)\Microsoft VS Code\Code.exe",
                os.path.expanduser(r"~\AppData\Local\Programs\Microsoft VS Code\Code.exe")
            ]
            for path in vscode_paths:
                if os.path.exists(path):
                    subprocess.Popen([path, str(filepath)])
                    return True
        
        # Fallback to default program
        try:
            if platform.system() == "Windows":
                os.startfile(str(filepath))
            elif platform.system() == "Darwin":
                os.system(f'open "{filepath}"')
            else:
                os.system(f'xdg-open "{filepath}"')
            return True
        except:
            return False

def run_generated_code(filename, language):
    """
    FIXED: Execute the generated code safely with proper error handling
    """
    try:
        speak(f"Running your {language} code now.")
        
        if language == "Python":
            # FIXED: Use absolute path and proper Python interpreter
            result = subprocess.run(
                ['python', str(filename)],
                capture_output=True,
                text=True,
                timeout=30,
                cwd=str(Path(filename).parent)  # Set working directory
            )
            
            if result.returncode == 0:
                speak("Code executed successfully!")
                if result.stdout:
                    print(f"\n=== OUTPUT ===\n{result.stdout}")
                    # Speak first 100 characters of output
                    output_preview = result.stdout[:100]
                    speak(f"Output: {output_preview}")
            else:
                speak("There was an error running the code.")
                print(f"\n=== ERROR ===\n{result.stderr}")
                # Speak the error
                error_msg = result.stderr.split('\n')[-2] if result.stderr else "Unknown error"
                speak(f"Error: {error_msg[:100]}")
                
        elif language == "JavaScript":
            try:
                result = subprocess.run(
                    ['node', str(filename)],
                    capture_output=True,
                    text=True,
                    timeout=30,
                    cwd=str(Path(filename).parent)
                )
                if result.returncode == 0:
                    speak("JavaScript code executed successfully!")
                    if result.stdout:
                        print(f"\n=== OUTPUT ===\n{result.stdout}")
                        speak(f"Output: {result.stdout[:100]}")
                else:
                    speak("Error running JavaScript code.")
                    print(f"\n=== ERROR ===\n{result.stderr}")
            except FileNotFoundError:
                speak("Node.js is not installed. Please install it to run JavaScript code.")
                print("Install Node.js from: https://nodejs.org/")
                
        else:
            speak(f"{language} code saved. Please compile and run it manually from VS Code.")
            
    except subprocess.TimeoutExpired:
        speak("Code execution timed out. It may be waiting for input or taking too long.")
    except Exception as e:
        print(f"Run error: {e}")
        import traceback
        traceback.print_exc()
        speak("Could not run the code. Check console for details.")




# News API Configuration (Optional)
NEWS_API_KEY = "20af6c9482b74c778676063604441009"  # Get from https://newsapi.org/register
NEWS_API_BASE_URL = "https://newsapi.org/v2/"
def get_news_hindi_from_rss(language="hindi"):
    """
    Fetch news using RSS feeds - Supports Hindi and English
    """
    try:
        import feedparser
        
        if language == "hindi":
            speak("हिंदी समाचार लाया जा रहा है")
            print("Fetching Hindi news from RSS feeds...")
            
            # Updated Hindi news RSS feeds with working URLs
            rss_feeds = {
                "आज तक": "https://www.aajtak.in/rss/top-stories",
                "BBC हिंदी": "https://feeds.bbci.co.uk/hindi/rss.xml",
                "न्यूज़18": "https://hindi.news18.com/rss/india.xml",
                "वन इंडिया": "https://hindi.oneindia.com/rss/news-fb.xml",
            }
        else:
            speak("Fetching latest news headlines")
            print("Fetching English news from RSS feeds...")
            
            # English news RSS feeds
            rss_feeds = {
                "BBC": "http://feeds.bbci.co.uk/news/rss.xml",
                "The Hindu": "https://www.thehindu.com/news/national/feeder/default.rss",
                "India Today": "https://www.indiatoday.in/rss/home",
                "NDTV": "https://feeds.feedburner.com/ndtvnews-top-stories",
            }
        
        all_articles = []
        
        for source_name, feed_url in rss_feeds.items():
            try:
                print(f"Fetching from {source_name}...")
                # Add timeout and user agent to avoid blocks
                feed = feedparser.parse(feed_url, 
                                       request_headers={'User-Agent': 'Mozilla/5.0'})
                
                # Check if feed was fetched successfully
                if feed.entries:
                    for entry in feed.entries[:3]:  # Get top 3 from each source
                        all_articles.append({
                            "source": source_name,
                            "title": entry.get("title", "No title"),
                            "summary": entry.get("summary", "")[:150]
                        })
                    print(f"✓ Got {len(feed.entries[:3])} articles from {source_name}")
                else:
                    print(f"⚠ No articles from {source_name}")
                    
            except Exception as e:
                print(f"✗ Error fetching from {source_name}: {e}")
                continue
        
        if all_articles:
            if language == "hindi":
                speak(f"यहां {len(all_articles)} मुख्य समाचार हैं:")
            else:
                speak(f"Here are the latest {len(all_articles)} headlines:")
            
            for i, article in enumerate(all_articles, 1):
                headline = f"{article['title']}"
                print(f"{i}. [{article['source']}] {headline}")
                
                # Speak in appropriate language
                if language == "hindi":
                    speak(f"समाचार {i}, {article['source']} से: {headline}", language="hindi")
                else:
                    speak(f"Headline {i} from {article['source']}: {headline}")
                time.sleep(1)
            
            return True
        else:
            if language == "hindi":
                speak("समाचार प्राप्त नहीं हो सका। कृपया अपना इंटरनेट कनेक्शन जांचें")
            else:
                speak("Could not fetch news. Please check your internet connection")
            return False
            
    except ImportError:
        print("\n⚠️ feedparser not installed")
        if language == "hindi":
            speak("कृपया feedparser लाइब्रेरी इंस्टॉल करें। कमांड प्रॉम्प्ट में टाइप करें: pip install feedparser")
        else:
            speak("Please install feedparser library. Type in command prompt: pip install feedparser")
        return False
    except Exception as e:
        print(f"RSS feed error: {e}")
        import traceback
        traceback.print_exc()
        if language == "hindi":
            speak("समाचार प्राप्त करने में त्रुटि हुई। इंटरनेट कनेक्शन जांचें")
        else:
            speak("Error fetching news. Please check your internet connection")
        return False
def get_news_from_rss():
    """
    Fetch news using RSS feeds - NO API KEY REQUIRED
    This will always work without any configuration
    """
    try:
        import feedparser
        
        speak("Fetching latest news headlines")
        print("Fetching news from RSS feeds...")
        
        # Popular Indian news RSS feeds
        rss_feeds = {
            "NDTV": "https://feeds.feedburner.com/ndtvnews-top-stories",
            "The Hindu": "https://www.thehindu.com/news/national/feeder/default.rss",
            "India Today": "https://www.indiatoday.in/rss/home",
        }
        
        all_articles = []
        
        for source_name, feed_url in rss_feeds.items():
            try:
                print(f"Fetching from {source_name}...")
                feed = feedparser.parse(feed_url)
                
                for entry in feed.entries[:3]:  # Get top 3 from each source
                    all_articles.append({
                        "source": source_name,
                        "title": entry.get("title", "No title"),
                        "summary": entry.get("summary", "")[:150]  # First 150 chars
                    })
            except Exception as e:
                print(f"Error fetching from {source_name}: {e}")
                continue
        
        if all_articles:
            speak(f"Here are the latest {len(all_articles)} headlines:")
            
            for i, article in enumerate(all_articles, 1):
                headline = f"{article['title']}"
                print(f"{i}. [{article['source']}] {headline}")
                speak(f"Headline {i} from {article['source']}: {headline}")
                time.sleep(1)
            
            return True
        else:
            speak("Could not fetch news from RSS feeds. Please check your internet connection")
            return False
            
    except ImportError:
        print("\n⚠️ feedparser not installed")
        speak("Please install feedparser library. Open command prompt and type: pip install feedparser")
        return False
    except Exception as e:
        print(f"RSS feed error: {e}")
        speak("Could not fetch news. Please check your internet connection")
        return False


def get_news_from_api(category="general", country="in", max_articles=5):
    """
    Fetch news using NewsAPI.org - REQUIRES API KEY
    Categories: business, entertainment, general, health, science, sports, technology
    """
    try:
        # Check if API key is configured
        if not NEWS_API_KEY or NEWS_API_KEY == "YOUR_NEWS_API_KEY_HERE":
            print("⚠️ NewsAPI key not configured, falling back to RSS feeds")
            return get_news_from_rss()
        
        speak(f"Fetching latest {category} news")
        print(f"Fetching {category} news from NewsAPI...")
        
        url = f"{NEWS_API_BASE_URL}top-headlines"
        params = {
            "country": country,
            "category": category,
            "apiKey": NEWS_API_KEY,
            "pageSize": max_articles
        }
        
        response = requests.get(url, params=params, timeout=15)
        
        if response.status_code == 200:
            data = response.json()
            articles = data.get("articles", [])
            
            if articles:
                speak(f"Here are the top {len(articles)} {category} headlines:")
                
                for i, article in enumerate(articles, 1):
                    title = article.get("title", "No title")
                    source = article.get("source", {}).get("name", "Unknown")
                    
                    print(f"{i}. [{source}] {title}")
                    speak(f"Headline {i}: {title}")
                    time.sleep(1)
                
                return True
            else:
                print(f"No articles found for category: {category}")
                speak(f"No {category} news found. Let me try general news instead")
                return get_news_from_rss()
        
        elif response.status_code == 401:
            print("❌ Invalid API key")
            speak("News API key is invalid. Switching to RSS feeds")
            return get_news_from_rss()
        
        elif response.status_code == 426:
            print("❌ API key requires upgrade")
            speak("News API requires upgrade. Using free RSS feeds instead")
            return get_news_from_rss()
        
        elif response.status_code == 429:
            print("❌ API rate limit exceeded")
            speak("News API limit reached. Using RSS feeds instead")
            return get_news_from_rss()
        
        else:
            print(f"API returned status code: {response.status_code}")
            speak("News API unavailable. Using RSS feeds instead")
            return get_news_from_rss()
            
    except requests.Timeout:
        print("⚠️ Request timed out")
        speak("Request timed out. Trying RSS feeds")
        return get_news_from_rss()
    except Exception as e:
        print(f"NewsAPI error: {e}")
        speak("News API error. Using RSS feeds instead")
        return get_news_from_rss()


def get_news_gnews_free():
    """
    Alternative free news source using web scraping techniques
    No API key needed
    """
    try:
        speak("Fetching latest news")
        
        # Using a public news aggregator
        url = "https://newsapi.org/v2/top-headlines?country=in&apiKey=demo"
        
        # Note: The demo key has very limited access
        response = requests.get(url, timeout=10)
        
        if response.status_code != 200:
            return get_news_from_rss()
        
        return True
        
    except Exception as e:
        print(f"Free news fetch error: {e}")
        return get_news_from_rss()


def search_news(query, max_articles=5):
    """
    Search for specific news topics
    """
    try:
        if not NEWS_API_KEY or NEWS_API_KEY == "YOUR_NEWS_API_KEY_HERE":
            speak(f"Searching for {query} in RSS feeds")
            return search_news_rss(query)
        
        speak(f"Searching news about {query}")
        
        url = f"{NEWS_API_BASE_URL}everything"
        params = {
            "q": query,
            "language": "en",
            "sortBy": "publishedAt",
            "apiKey": NEWS_API_KEY,
            "pageSize": max_articles
        }
        
        response = requests.get(url, params=params, timeout=15)
        
        if response.status_code == 200:
            data = response.json()
            articles = data.get("articles", [])
            
            if articles:
                speak(f"Found {len(articles)} articles about {query}:")
                
                for i, article in enumerate(articles, 1):
                    title = article.get("title", "No title")
                    source = article.get("source", {}).get("name", "Unknown")
                    
                    print(f"{i}. [{source}] {title}")
                    speak(f"Article {i}: {title}")
                    time.sleep(1)
                
                return True
            else:
                speak(f"No articles found about {query}")
                return False
        else:
            return search_news_rss(query)
            
    except Exception as e:
        print(f"News search error: {e}")
        return search_news_rss(query)


def search_news_rss(query):
    """
    Search news in RSS feeds (basic filtering)
    """
    try:
        import feedparser
        
        speak(f"Searching for {query} in news feeds")
        
        feeds = [
            "https://feeds.feedburner.com/ndtvnews-top-stories",
            "https://www.thehindu.com/news/national/feeder/default.rss",
        ]
        
        matching_articles = []
        query_lower = query.lower()
        
        for feed_url in feeds:
            try:
                feed = feedparser.parse(feed_url)
                for entry in feed.entries:
                    title = entry.get("title", "").lower()
                    if query_lower in title:
                        matching_articles.append(entry.get("title", ""))
            except:
                continue
        
        if matching_articles:
            speak(f"Found {len(matching_articles)} articles about {query}")
            for i, title in enumerate(matching_articles[:5], 1):
                speak(f"Article {i}: {title}")
                time.sleep(1)
            return True
        else:
            speak(f"No articles found about {query}")
            return False
            
    except ImportError:
        speak("Cannot search news. Please install feedparser")
        return False
    except Exception as e:
        print(f"RSS search error: {e}")
        return False


def get_current_affairs():
    """
    Main function - automatically chooses best method
    ALWAYS WORKS - uses RSS feeds as reliable fallback
    """
    try:
        # Try API first if configured
        if NEWS_API_KEY and NEWS_API_KEY != "YOUR_NEWS_API_KEY_HERE":
            print("Trying NewsAPI...")
            return get_news_from_api(category="general", country="in")
        else:
            # Use RSS feeds (always works)
            print("Using RSS feeds (no API key configured)...")
            return get_news_from_rss()
            
    except Exception as e:
        print(f"Current affairs error: {e}")
        # Last resort - always try RSS
        return get_news_from_rss()


def get_business_news():
    """Get business news"""
    return get_news_from_api(category="business", country="in")


def get_sports_news():
    """Get sports news"""
    return get_news_from_api(category="sports", country="in")


def get_tech_news():
    """Get technology news"""
    return get_news_from_api(category="technology", country="in")
MEDICAL_CONDITIONS = {
    "fever": {
        "en": {
            "symptoms": "High body temperature, sweating, chills, headache, muscle aches",
            "causes": "Infection, heat exhaustion, inflammatory conditions",
            "first_aid": "Rest, drink plenty of fluids, take fever-reducing medication like paracetamol",
            "when_to_see_doctor": "If fever is above 103°F, lasts more than 3 days, or accompanied by severe symptoms"
        },
        "hi": {
            "symptoms": "उच्च शरीर का तापमान, पसीना आना, ठंड लगना, सिरदर्द, मांसपेशियों में दर्द",
            "causes": "संक्रमण, गर्मी से थकावट, सूजन संबंधी स्थितियां",
            "first_aid": "आराम करें, खूब तरल पदार्थ पिएं, पेरासिटामोल जैसी बुखार कम करने वाली दवा लें",
            "when_to_see_doctor": "यदि बुखार 103°F से अधिक है, 3 दिन से अधिक समय तक रहता है, या गंभीर लक्षणों के साथ है"
        }
    },
    "headache": {
        "en": {
            "symptoms": "Pain in head, sensitivity to light, nausea",
            "causes": "Stress, dehydration, eye strain, lack of sleep, tension",
            "first_aid": "Rest in a quiet dark room, drink water, apply cold compress, take pain reliever",
            "when_to_see_doctor": "If sudden severe headache, headache with fever, confusion, or vision changes"
        },
        "hi": {
            "symptoms": "सिर में दर्द, प्रकाश के प्रति संवेदनशीलता, मतली",
            "causes": "तनाव, निर्जलीकरण, आंखों में खिंचाव, नींद की कमी",
            "first_aid": "शांत अंधेरे कमरे में आराम करें, पानी पिएं, ठंडी पट्टी लगाएं, दर्द निवारक लें",
            "when_to_see_doctor": "यदि अचानक गंभीर सिरदर्द, बुखार के साथ सिरदर्द, भ्रम, या दृष्टि में बदलाव"
        }
    },
    "cold": {
        "en": {
            "symptoms": "Runny nose, sneezing, sore throat, cough, mild fever",
            "causes": "Viral infection, weak immunity, exposure to cold weather",
            "first_aid": "Rest, drink warm fluids, gargle with salt water, take vitamin C",
            "when_to_see_doctor": "If symptoms persist beyond 10 days or worsen significantly"
        },
        "hi": {
            "symptoms": "नाक बहना, छींक आना, गले में खराश, खांसी, हल्का बुखार",
            "causes": "वायरल संक्रमण, कमजोर प्रतिरक्षा, ठंडे मौसम के संपर्क में आना",
            "first_aid": "आराम करें, गर्म तरल पदार्थ पिएं, नमक के पानी से गरारे करें, विटामिन सी लें",
            "when_to_see_doctor": "यदि लक्षण 10 दिनों से अधिक समय तक बने रहें या काफी खराब हो जाएं"
        }
    },
    "stomach pain": {
        "en": {
            "symptoms": "Abdominal discomfort, cramping, nausea, bloating",
            "causes": "Indigestion, gas, food poisoning, stress, infection",
            "first_aid": "Rest, avoid solid food initially, drink clear fluids, apply warm compress",
            "when_to_see_doctor": "If severe pain, blood in stool, persistent vomiting, or pain lasts more than 24 hours"
        },
        "hi": {
            "symptoms": "पेट में बेचैनी, ऐंठन, मतली, सूजन",
            "causes": "अपच, गैस, खाद्य विषाक्तता, तनाव, संक्रमण",
            "first_aid": "आराम करें, शुरू में ठोस भोजन से बचें, साफ तरल पदार्थ पिएं, गर्म सेक लगाएं",
            "when_to_see_doctor": "यदि गंभीर दर्द, मल में रक्त, लगातार उल्टी, या दर्द 24 घंटे से अधिक समय तक रहे"
        }
    },
    "diabetes": {
        "en": {
            "symptoms": "Increased thirst, frequent urination, fatigue, blurred vision, slow healing wounds",
            "causes": "Insulin resistance, pancreatic problems, genetic factors, lifestyle",
            "management": "Regular monitoring, balanced diet, exercise, medication as prescribed",
            "when_to_see_doctor": "Regular check-ups, if blood sugar consistently high or low"
        },
        "hi": {
            "symptoms": "अत्यधिक प्यास, बार-बार पेशाब आना, थकान, धुंधली दृष्टि, घाव धीरे-धीरे भरना",
            "causes": "इंसुलिन प्रतिरोध, अग्न्याशय की समस्याएं, आनुवंशिक कारक, जीवन शैली",
            "management": "नियमित निगरानी, संतुलित आहार, व्यायाम, निर्धारित दवाएं",
            "when_to_see_doctor": "नियमित जांच, यदि रक्त शर्करा लगातार उच्च या निम्न"
        }
    }
}

# Legal Rights Database
LEGAL_RIGHTS = {
    "right to equality": {
        "en": {
            "article": "Article 14-18 of Indian Constitution",
            "description": "Equality before law, prohibition of discrimination on grounds of religion, race, caste, sex or place of birth",
            "what_it_means": "Every citizen is equal before the law and has equal protection of laws",
            "violations": "Discrimination in employment, education, or public services based on caste, religion, or gender"
        },
        "hi": {
            "article": "भारतीय संविधान के अनुच्छेद 14-18",
            "description": "कानून के समक्ष समानता, धर्म, जाति, लिंग या जन्म स्थान के आधार पर भेदभाव का निषेध",
            "what_it_means": "प्रत्येक नागरिक कानून के समक्ष समान है और कानूनों की समान सुरक्षा प्राप्त है",
            "violations": "जाति, धर्म या लिंग के आधार पर रोजगार, शिक्षा या सार्वजनिक सेवाओं में भेदभाव"
        }
    },
    "right to freedom": {
        "en": {
            "article": "Article 19-22 of Indian Constitution",
            "description": "Freedom of speech, assembly, association, movement, residence, and profession",
            "what_it_means": "Citizens have freedom to express opinions, form associations, move freely, and choose occupation",
            "limitations": "Subject to reasonable restrictions for security, public order, and morality"
        },
        "hi": {
            "article": "भारतीय संविधान के अनुच्छेद 19-22",
            "description": "भाषण, सभा, संघ बनाने, आवागमन, निवास और पेशे की स्वतंत्रता",
            "what_it_means": "नागरिकों को विचार व्यक्त करने, संघ बनाने, स्वतंत्र रूप से घूमने और पेशा चुनने की स्वतंत्रता है",
            "limitations": "सुरक्षा, सार्वजनिक व्यवस्था और नैतिकता के लिए उचित प्रतिबंधों के अधीन"
        }
    },
    "right against exploitation": {
        "en": {
            "article": "Article 23-24 of Indian Constitution",
            "description": "Prohibition of human trafficking, forced labor, and child labor",
            "what_it_means": "No person can be forced to work without payment or in hazardous conditions",
            "violations": "Bonded labor, child labor in factories, human trafficking"
        },
        "hi": {
            "article": "भारतीय संविधान के अनुच्छेद 23-24",
            "description": "मानव तस्करी, बलात्कार श्रम और बाल श्रम का निषेध",
            "what_it_means": "किसी भी व्यक्ति को बिना भुगतान या खतरनाक परिस्थितियों में काम करने के लिए मजबूर नहीं किया जा सकता",
            "violations": "बंधुआ मजदूरी, कारखानों में बाल श्रम, मानव तस्करी"
        }
    },
    "consumer rights": {
        "en": {
            "description": "Rights to safety, information, choice, and redressal",
            "what_it_means": "Protection against defective products, unfair trade practices, right to compensation",
            "how_to_complain": "File complaint with Consumer Forum, National Consumer Helpline: 1800-11-4000",
            "time_limit": "Complaint within 2 years of purchase or issue"
        },
        "hi": {
            "description": "सुरक्षा, सूचना, चयन और निवारण का अधिकार",
            "what_it_means": "दोषपूर्ण उत्पादों, अनुचित व्यापार प्रथाओं के खिलाफ सुरक्षा, मुआवजे का अधिकार",
            "how_to_complain": "उपभोक्ता फोरम में शिकायत दर्ज करें, राष्ट्रीय उपभोक्ता हेल्पलाइन: 1800-11-4000",
            "time_limit": "खरीद या समस्या के 2 साल के भीतर शिकायत करें"
        }
    },
    "women rights": {
        "en": {
            "description": "Protection against domestic violence, workplace harassment, and discrimination",
            "laws": "Domestic Violence Act 2005, Sexual Harassment Act 2013, Dowry Prohibition Act",
            "helpline": "Women Helpline: 1091, National Commission for Women: 011-26942369",
            "what_to_do": "File FIR at nearest police station, contact women helpline, seek legal aid"
        },
        "hi": {
            "description": "घरेलू हिंसा, कार्यस्थल उत्पीड़न और भेदभाव के खिलाफ सुरक्षा",
            "laws": "घरेलू हिंसा अधिनियम 2005, यौन उत्पीड़न अधिनियम 2013, दहेज निषेध अधिनियम",
            "helpline": "महिला हेल्पलाइन: 1091, राष्ट्रीय महिला आयोग: 011-26942369",
            "what_to_do": "निकटतम पुलिस स्टेशन में एफआईआर दर्ज करें, महिला हेल्पलाइन से संपर्क करें, कानूनी सहायता लें"
        }
    }
}

# Emergency Contact Numbers
EMERGENCY_CONTACTS = {
    "en": {
        "police": "100",
        "ambulance": "102 or 108",
        "fire": "101",
        "women_helpline": "1091",
        "child_helpline": "1098",
        "senior_citizen": "1091 or 1291",
        "disaster_management": "108",
        "cyber_crime": "1930",
        "national_helpline": "112"
    },
    "hi": {
        "police": "100 - पुलिस",
        "ambulance": "102 या 108 - एम्बुलेंस",
        "fire": "101 - अग्निशमन",
        "women_helpline": "1091 - महिला हेल्पलाइन",
        "child_helpline": "1098 - बाल हेल्पलाइन",
        "senior_citizen": "1091 या 1291 - वरिष्ठ नागरिक",
        "disaster_management": "108 - आपदा प्रबंधन",
        "cyber_crime": "1930 - साइबर अपराध",
        "national_helpline": "112 - राष्ट्रीय हेल्पलाइन"
    }
}

def get_medical_info(condition, language="en"):
    """
    Get medical information about a condition
    """
    try:
        condition = condition.lower().strip()
        
        # Search in database
        for key in MEDICAL_CONDITIONS.keys():
            if key in condition or condition in key:
                lang_code = "hi" if language == "hindi" else "en"
                info = MEDICAL_CONDITIONS[key].get(lang_code, MEDICAL_CONDITIONS[key]["en"])
                
                if language == "hindi":
                    speak(f"{key} के बारे में जानकारी:", language="hindi")
                    time.sleep(0.3)
                    
                    if "symptoms" in info:
                        speak(f"लक्षण: {info['symptoms']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "causes" in info:
                        speak(f"कारण: {info['causes']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "first_aid" in info:
                        speak(f"प्राथमिक उपचार: {info['first_aid']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "management" in info:
                        speak(f"प्रबंधन: {info['management']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "when_to_see_doctor" in info:
                        speak(f"डॉक्टर से कब मिलें: {info['when_to_see_doctor']}", language="hindi")
                else:
                    speak(f"Information about {key}:")
                    time.sleep(0.3)
                    
                    if "symptoms" in info:
                        speak(f"Symptoms: {info['symptoms']}")
                        time.sleep(0.3)
                    
                    if "causes" in info:
                        speak(f"Causes: {info['causes']}")
                        time.sleep(0.3)
                    
                    if "first_aid" in info:
                        speak(f"First aid: {info['first_aid']}")
                        time.sleep(0.3)
                    
                    if "management" in info:
                        speak(f"Management: {info['management']}")
                        time.sleep(0.3)
                    
                    if "when_to_see_doctor" in info:
                        speak(f"When to see doctor: {info['when_to_see_doctor']}")
                
                return True
        
        # If not found, use AI to generate information
        if language == "hindi":
            speak(f"{condition} के बारे में जानकारी खोजी जा रही है", language="hindi")
        else:
            speak(f"Searching for information about {condition}")
        
        return search_medical_info_online(condition, language)
        
    except Exception as e:
        print(f"Medical info error: {e}")
        if language == "hindi":
            speak("चिकित्सा जानकारी प्राप्त नहीं हो सकी", language="hindi")
        else:
            speak("Could not retrieve medical information")
        return False


def search_medical_info_online(condition, language="en"):
    """
    Search medical information online using AI
    """
    try:
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        
        if language == "hindi":
            prompt = f""""{condition}" के बारे में संक्षिप्त चिकित्सा जानकारी प्रदान करें:

1. मुख्य लक्षण (2-3 लाइन में)
2. सामान्य कारण (2-3 लाइन में)
3. घरेलू उपचार / प्राथमिक उपचार (2-3 लाइन में)
4. डॉक्टर से कब मिलना चाहिए (1-2 लाइन में)

सरल हिंदी में जवाब दें। कुल 200 शब्दों से कम में।

महत्वपूर्ण: यह केवल सामान्य जानकारी है। गंभीर स्थिति में तुरंत डॉक्टर से संपर्क करें।"""
        else:
            prompt = f"""Provide brief medical information about "{condition}":

1. Main symptoms (2-3 lines)
2. Common causes (2-3 lines)
3. Home remedies / First aid (2-3 lines)
4. When to see a doctor (1-2 lines)

Keep response under 200 words and in simple language.

IMPORTANT: This is general information only. For serious conditions, consult a doctor immediately."""

        response = model.generate_content(prompt)
        info_text = response.text.strip()
        
        # Speak in chunks
        sentences = info_text.split('.')
        for sentence in sentences:
            if sentence.strip():
                if language == "hindi":
                    speak(sentence.strip(), language="hindi")
                else:
                    speak(sentence.strip())
                time.sleep(0.3)
        
        return True
        
    except Exception as e:
        print(f"Online medical search error: {e}")
        return False


def get_legal_info(right_or_law, language="en"):
    """
    Get information about legal rights and laws
    """
    try:
        query = right_or_law.lower().strip()
        
        # Search in database
        for key in LEGAL_RIGHTS.keys():
            if key in query or any(word in query for word in key.split()):
                lang_code = "hi" if language == "hindi" else "en"
                info = LEGAL_RIGHTS[key].get(lang_code, LEGAL_RIGHTS[key]["en"])
                
                if language == "hindi":
                    speak(f"{key} के बारे में जानकारी:", language="hindi")
                    time.sleep(0.3)
                    
                    if "article" in info:
                        speak(f"अनुच्छेद: {info['article']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "description" in info:
                        speak(f"विवरण: {info['description']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "what_it_means" in info:
                        speak(f"इसका मतलब: {info['what_it_means']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "laws" in info:
                        speak(f"कानून: {info['laws']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "helpline" in info:
                        speak(f"हेल्पलाइन: {info['helpline']}", language="hindi")
                        time.sleep(0.3)
                    
                    if "violations" in info:
                        speak(f"उल्लंघन: {info['violations']}", language="hindi")
                    
                    if "how_to_complain" in info:
                        speak(f"शिकायत कैसे करें: {info['how_to_complain']}", language="hindi")
                else:
                    speak(f"Information about {key}:")
                    time.sleep(0.3)
                    
                    if "article" in info:
                        speak(f"Article: {info['article']}")
                        time.sleep(0.3)
                    
                    if "description" in info:
                        speak(f"Description: {info['description']}")
                        time.sleep(0.3)
                    
                    if "what_it_means" in info:
                        speak(f"What it means: {info['what_it_means']}")
                        time.sleep(0.3)
                    
                    if "laws" in info:
                        speak(f"Laws: {info['laws']}")
                        time.sleep(0.3)
                    
                    if "helpline" in info:
                        speak(f"Helpline: {info['helpline']}")
                        time.sleep(0.3)
                    
                    if "violations" in info:
                        speak(f"Violations: {info['violations']}")
                    
                    if "how_to_complain" in info:
                        speak(f"How to complain: {info['how_to_complain']}")
                
                return True
        
        # If not found, search online
        if language == "hindi":
            speak(f"{right_or_law} के बारे में खोज रहे हैं", language="hindi")
        else:
            speak(f"Searching for information about {right_or_law}")
        
        return search_legal_info_online(right_or_law, language)
        
    except Exception as e:
        print(f"Legal info error: {e}")
        if language == "hindi":
            speak("कानूनी जानकारी प्राप्त नहीं हो सकी", language="hindi")
        else:
            speak("Could not retrieve legal information")
        return False


def search_legal_info_online(query, language="en"):
    """
    Search legal information online using AI
    """
    try:
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        
        if language == "hindi":
            prompt = f"""भारतीय कानून के अनुसार "{query}" के बारे में संक्षिप्त जानकारी दें:

1. यह अधिकार/कानून क्या है? (2-3 लाइन)
2. संबंधित अनुच्छेद/कानून (यदि लागू हो)
3. इसका मतलब सरल शब्दों में (2-3 लाइन)
4. संबंधित हेल्पलाइन या शिकायत कैसे करें (यदि लागू हो)

कुल 200 शब्दों से कम में उत्तर दें।

नोट: यह सामान्य जानकारी है। कानूनी सलाह के लिए वकील से परामर्श करें।"""
        else:
            prompt = f"""Provide brief information about "{query}" according to Indian law:

1. What is this right/law? (2-3 lines)
2. Related articles/acts (if applicable)
3. What it means in simple terms (2-3 lines)
4. Related helplines or how to complain (if applicable)

Keep response under 200 words.

Note: This is general information. Consult a lawyer for legal advice."""

        response = model.generate_content(prompt)
        info_text = response.text.strip()
        
        # Speak in chunks
        sentences = info_text.split('.')
        for sentence in sentences:
            if sentence.strip():
                if language == "hindi":
                    speak(sentence.strip(), language="hindi")
                else:
                    speak(sentence.strip())
                time.sleep(0.3)
        
        return True
        
    except Exception as e:
        print(f"Online legal search error: {e}")
        return False


def list_emergency_numbers(language="en"):
    """
    List all emergency helpline numbers
    """
    try:
        lang_code = "hi" if language == "hindi" else "en"
        numbers = EMERGENCY_CONTACTS[lang_code]
        
        if language == "hindi":
            speak("आपातकालीन नंबर:", language="hindi")
            time.sleep(0.3)
            
            for service, number in numbers.items():
                speak(number, language="hindi")
                time.sleep(0.5)
        else:
            speak("Emergency helpline numbers:")
            time.sleep(0.3)
            
            speak(f"Police: {numbers['police']}")
            time.sleep(0.3)
            speak(f"Ambulance: {numbers['ambulance']}")
            time.sleep(0.3)
            speak(f"Fire: {numbers['fire']}")
            time.sleep(0.3)
            speak(f"Women Helpline: {numbers['women_helpline']}")
            time.sleep(0.3)
            speak(f"Child Helpline: {numbers['child_helpline']}")
            time.sleep(0.3)
            speak(f"National Emergency Number: {numbers['national_helpline']}")
        
        return True
        
    except Exception as e:
        print(f"Emergency numbers error: {e}")
        speak("Could not retrieve emergency numbers")
        return False


def get_first_aid_info(emergency_type, language="en"):
    """
    Get first aid instructions for emergencies
    """
    try:
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        
        if language == "hindi":
            prompt = f""""{emergency_type}" के लिए प्राथमिक चिकित्सा निर्देश दें:

1. तुरंत क्या करें (3-4 स्टेप्स)
2. क्या नहीं करना चाहिए (2-3 पॉइंट्स)
3. कब तुरंत डॉक्टर बुलाएं

सरल हिंदी में, स्पष्ट स्टेप्स में लिखें। कुल 150 शब्दों में।

चेतावनी: यह सामान्य जानकारी है। गंभीर स्थिति में तुरंत 102 या 108 पर कॉल करें।"""
        else:
            prompt = f"""Provide first aid instructions for "{emergency_type}":

1. What to do immediately (3-4 steps)
2. What NOT to do (2-3 points)
3. When to call emergency services immediately

Write in simple language with clear steps. Keep under 150 words.

WARNING: This is general information. In serious emergencies, call 102 or 108 immediately."""

        response = model.generate_content(prompt)
        instructions = response.text.strip()
        
        # Speak in sections
        sections = instructions.split('\n')
        for section in sections:
            if section.strip():
                if language == "hindi":
                    speak(section.strip(), language="hindi")
                else:
                    speak(section.strip())
                time.sleep(0.4)
        
        return True
        
    except Exception as e:
        print(f"First aid info error: {e}")
        if language == "hindi":
            speak("प्राथमिक चिकित्सा जानकारी प्राप्त नहीं हो सकी", language="hindi")
        else:
            speak("Could not retrieve first aid information")
        return False


def get_entertainment_news():
    """Get entertainment news"""
    return get_news_from_api(category="entertainment", country="in")


import speech_recognition as sr
from pathlib import Path
import datetime
import time
import subprocess
import os
import pyautogui

def find_word_path():
    """
    Find Microsoft Word installation path
    """
    possible_paths = [
        r"C:\Program Files\Microsoft Office\root\Office16\WINWORD.EXE",
        r"C:\Program Files (x86)\Microsoft Office\root\Office16\WINWORD.EXE",
        r"C:\Program Files\Microsoft Office\Office16\WINWORD.EXE",
        r"C:\Program Files (x86)\Microsoft Office\Office16\WINWORD.EXE",
        r"C:\Program Files\Microsoft Office\root\Office15\WINWORD.EXE",
        r"C:\Program Files (x86)\Microsoft Office\root\Office15\WINWORD.EXE",
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            print(f"Found Word at: {path}")
            return path
    
    return None


def open_word_and_type(text):
    """
    Type text directly using keyboard automation
    """
    try:
        time.sleep(0.5)
        pyautogui.write(text + " ", interval=0.01)
        return True
    except Exception as e:
        print(f"Typing error: {e}")
        return False


def dictate_to_word_live():
    """
    Open Microsoft Word and write in real-time using keyboard automation
    Works with any Word version
    """
    try:
        speak("Opening Microsoft Word for live dictation")
        
        # Try to find and open Word
        word_path = find_word_path()
        
        if word_path:
            # Open Word using path
            subprocess.Popen([word_path])
            print("Word opened using direct path")
        else:
            # Try opening with system command
            try:
                os.system("start winword")
                print("Word opened using system command")
            except:
                speak("Could not find Microsoft Word. Please open Word manually and say 'ready'")
                return False
        
        speak("Waiting for Word to open. This will take a few seconds.")
        time.sleep(5)  # Wait for Word to fully open
        
        speak("Click on the Word document where you want to type, then say 'start'")
        print("\n" + "="*60)
        print("PREPARATION:")
        print("1. Word should now be open")
        print("2. Click inside the document")
        print("3. Say 'start' to begin dictation")
        print("="*60)
        
        # Wait for user to say "start"
        r = sr.Recognizer()
        ready = False
        
        for _ in range(5):  # Try 5 times
            try:
                with sr.Microphone() as source:
                    print("Say 'start' when ready...")
                    r.adjust_for_ambient_noise(source)
                    audio = r.listen(source, timeout=10)
                    command = r.recognize_google(audio, language=CURRENT_LANG_CODE)
                    if "start" in command.lower():
                        ready = True
                        break
            except:
                continue
        
        if not ready:
            speak("Starting anyway. Make sure Word document is active.")
            time.sleep(2)
        
        speak("Starting dictation. Say 'stop dictation' when finished.")
        print("\n" + "="*60)
        print("LIVE DICTATION ACTIVE")
        print("Commands:")
        print("  - Say 'new line' for line break")
        print("  - Say 'new paragraph' for paragraph")
        print("  - Say 'stop dictation' to end")
        print("="*60 + "\n")
        
        sentence_count = 0
        
        while True:
            try:
                with sr.Microphone() as source:
                    print(f"Listening... (Sentences: {sentence_count})")
                    eel.DisplayMessage('Listening and typing...')
                    r.adjust_for_ambient_noise(source, duration=0.5)
                    audio = r.listen(source, timeout=10, phrase_time_limit=15)
                
                print("Processing...")
                text = r.recognize_google(audio, language=CURRENT_LANG_CODE)
                text_lower = text.lower()
                
                print(f"You said: {text}")
                eel.DisplayMessage(f"Typing: {text}")
                
                # Check for stop command
                if "stop dictation" in text_lower or "end dictation" in text_lower:
                    speak("Stopping dictation")
                    break
                
                # Check for formatting commands
                if "new paragraph" in text_lower:
                    pyautogui.press('enter')
                    pyautogui.press('enter')
                    speak("New paragraph")
                    continue
                
                elif "new line" in text_lower:
                    pyautogui.press('enter')
                    speak("New line")
                    continue
                
                # Type the text
                open_word_and_type(text)
                sentence_count += 1
                
                time.sleep(0.3)
                
            except sr.WaitTimeoutError:
                print("No speech detected...")
                continue
                
            except sr.UnknownValueError:
                print("Could not understand")
                continue
                
            except Exception as e:
                print(f"Error: {e}")
                continue
        
        speak(f"Dictation complete. Typed {sentence_count} sentences. You can save the document now.")
        print(f"\nTotal sentences typed: {sentence_count}")
        return True
            
    except Exception as e:
        print(f"Dictation error: {e}")
        import traceback
        traceback.print_exc()
        speak("Sorry, I encountered an error during dictation")
        return False


def dictate_with_punctuation():
    """
    Enhanced dictation with voice punctuation
    """
    try:
        speak("Opening Word for dictation with punctuation support")
        
        word_path = find_word_path()
        if word_path:
            subprocess.Popen([word_path])
        else:
            os.system("start winword")
        
        speak("Waiting for Word to open")
        time.sleep(5)
        
        speak("Click in the document, then say 'ready'")
        
        r = sr.Recognizer()
        
        # Wait for ready
        for _ in range(5):
            try:
                with sr.Microphone() as source:
                    print("Say 'ready' when you're set...")
                    r.adjust_for_ambient_noise(source)
                    audio = r.listen(source, timeout=10)
                    command = r.recognize_google(audio, language=CURRENT_LANG_CODE)
                    if "ready" in command.lower():
                        break
            except:
                continue
        
        speak("Starting dictation with punctuation. Say 'comma', 'period', 'question mark', etc.")
        print("\n" + "="*60)
        print("PUNCTUATION DICTATION")
        print("Say: comma, period, question mark, exclamation mark")
        print("     new line, new paragraph, stop dictation")
        print("="*60 + "\n")
        
        word_count = 0
        
        punctuation_map = {
            "comma": ",",
            "period": ".",
            "full stop": ".",
            "question mark": "?",
            "exclamation mark": "!",
            "exclamation": "!",
            "colon": ":",
            "semicolon": ";"
        }
        
        while True:
            try:
                with sr.Microphone() as source:
                    print(f"Listening... (Words: {word_count})")
                    r.adjust_for_ambient_noise(source)
                    audio = r.listen(source, timeout=10, phrase_time_limit=15)
                
                text = r.recognize_google(audio, language=CURRENT_LANG_CODE)
                text_lower = text.lower()
                
                print(f"Recognized: {text}")
                
                if "stop dictation" in text_lower:
                    speak("Ending dictation")
                    break
                
                # Formatting
                if "new paragraph" in text_lower:
                    pyautogui.press('enter')
                    pyautogui.press('enter')
                    continue
                
                elif "new line" in text_lower:
                    pyautogui.press('enter')
                    continue
                
                # Check punctuation
                punctuation_found = False
                for phrase, symbol in punctuation_map.items():
                    if phrase in text_lower:
                        # Remove punctuation word
                        clean_text = text.lower().replace(phrase, "").strip()
                        if clean_text:
                            pyautogui.write(clean_text + symbol + " ", interval=0.01)
                            word_count += len(clean_text.split())
                        else:
                            pyautogui.write(symbol + " ", interval=0.01)
                        punctuation_found = True
                        break
                
                if not punctuation_found:
                    pyautogui.write(text + " ", interval=0.01)
                    word_count += len(text.split())
                
                time.sleep(0.3)
                
            except sr.WaitTimeoutError:
                continue
            except sr.UnknownValueError:
                continue
            except Exception as e:
                print(f"Error: {e}")
                continue
        
        speak(f"Dictation complete with {word_count} words")
        return True
        
    except Exception as e:
        print(f"Punctuation dictation error: {e}")
        speak("Error during dictation")
        return False


def simple_dictation():
    """
    Simplest version - just type what you say
    User must have Word open already
    """
    try:
        speak("Please open Microsoft Word and click in the document")
        speak("Say 'ready' when you want to start")
        
        r = sr.Recognizer()
        
        # Wait for ready
        ready = False
        for _ in range(10):
            try:
                with sr.Microphone() as source:
                    print("Waiting for 'ready' command...")
                    r.adjust_for_ambient_noise(source)
                    audio = r.listen(source, timeout=8)
                    command = r.recognize_google(audio, language=CURRENT_LANG_CODE)
                    if "ready" in command.lower() or "start" in command.lower():
                        ready = True
                        break
            except:
                continue
        
        if not ready:
            speak("Timeout. Please try again.")
            return False
        
        speak("Starting simple dictation. Say stop when done.")
        print("\n" + "="*60)
        print("SIMPLE DICTATION MODE")
        print("Say 'stop' or 'stop dictation' to end")
        print("="*60 + "\n")
        
        count = 0
        
        while True:
            try:
                with sr.Microphone() as source:
                    print(f"Listening... ({count} phrases)")
                    r.adjust_for_ambient_noise(source)
                    audio = r.listen(source, timeout=10, phrase_time_limit=15)
                
                text = r.recognize_google(audio, language=CURRENT_LANG_CODE)
                
                if "stop" in text.lower():
                    speak("Stopping")
                    break
                
                print(f"Typing: {text}")
                pyautogui.write(text + " ", interval=0.01)
                count += 1
                
                time.sleep(0.2)
                
            except sr.WaitTimeoutError:
                continue
            except sr.UnknownValueError:
                continue
            except Exception as e:
                print(f"Error: {e}")
                continue
        
        speak(f"Typed {count} phrases")
        return True
        
    except Exception as e:
        print(f"Simple dictation error: {e}")
        speak("Error during dictation")
        return False

def get_system_info():
    """
    Get comprehensive system information
    """
    try:
        speak("Gathering system information")
        
        # Operating System Info
        os_name = platform.system()
        os_version = platform.version()
        os_release = platform.release()
        
        # Computer Info
        computer_name = platform.node()
        processor = platform.processor()
        architecture = platform.machine()
        
        # Python Info
        python_version = platform.python_version()
        
        # Memory Info
        memory = psutil.virtual_memory()
        total_ram = round(memory.total / (1024**3), 2)
        available_ram = round(memory.available / (1024**3), 2)
        ram_percent = memory.percent
        
        # Disk Info
        disk = psutil.disk_usage('C:\\' if os_name == 'Windows' else '/')
        total_disk = round(disk.total / (1024**3), 2)
        free_disk = round(disk.free / (1024**3), 2)
        disk_percent = disk.percent
        
        # CPU Info
        cpu_count = psutil.cpu_count()
        cpu_percent = psutil.cpu_percent(interval=1)
        
        info_text = f"""
System Information:

Operating System: {os_name} {os_release}
Version: {os_version}
Computer Name: {computer_name}
Processor: {processor}
Architecture: {architecture}
CPU Cores: {cpu_count}
CPU Usage: {cpu_percent}%

Python Version: {python_version}

RAM: {total_ram} GB Total, {available_ram} GB Available ({ram_percent}% used)
Disk C: {total_disk} GB Total, {free_disk} GB Free ({disk_percent}% used)
"""
        
        print(info_text)
        
        # Speak summary
        speak(f"System: {os_name} {os_release}")
        speak(f"Processor: {processor}")
        speak(f"RAM: {total_ram} gigabytes total, {ram_percent} percent used")
        speak(f"Disk: {free_disk} gigabytes free out of {total_disk} gigabytes")
        speak(f"CPU usage: {cpu_percent} percent")
        
        return True
        
    except Exception as e:
        print(f"System info error: {e}")
        speak("Could not retrieve complete system information")
        return False


# ==================== WINDOWS UPDATE ====================

def update_windows():
    """
    Check and install Windows updates
    Requires administrator privileges
    """
    try:
        if platform.system() != "Windows":
            speak("Windows update is only available on Windows operating system")
            return False
        
        speak("Checking for Windows updates. This may take a few minutes.")
        print("\n" + "="*60)
        print("WINDOWS UPDATE")
        print("="*60)
        
        # Method 1: Using PowerShell (Recommended)
        powershell_cmd = """
        Install-Module PSWindowsUpdate -Force -Scope CurrentUser
        Import-Module PSWindowsUpdate
        Get-WindowsUpdate -Install -AcceptAll -AutoReboot
        """
        
        try:
            # Check for updates first
            check_cmd = "Get-WindowsUpdate"
            result = subprocess.run(
                ["powershell", "-Command", check_cmd],
                capture_output=True,
                text=True,
                timeout=60
            )
            
            if "No updates available" in result.stdout or not result.stdout.strip():
                speak("Your Windows is up to date. No updates found.")
                print("✓ Windows is up to date")
                return True
            
            speak("Updates available. Would you like to install them? Say yes or no.")
            confirmation = takecommand()
            
            if "yes" in confirmation.lower():
                speak("Installing Windows updates. This will take some time. Please be patient.")
                
                # Install updates
                result = subprocess.run(
                    ["powershell", "-Command", powershell_cmd],
                    capture_output=True,
                    text=True,
                    timeout=1800  # 30 minutes timeout
                )
                
                if result.returncode == 0:
                    speak("Windows updates installed successfully. A restart may be required.")
                    print("✓ Windows updates completed")
                    return True
                else:
                    speak("Some updates may require administrator privileges")
                    print(f"⚠️ Update result: {result.stderr}")
            else:
                speak("Windows update cancelled")
                return False
                
        except subprocess.TimeoutExpired:
            speak("Update process is taking longer than expected. Continuing in background.")
            return True
        except Exception as e:
            print(f"PowerShell update error: {e}")
            
            # Method 2: Alternative using UsoClient (Windows built-in)
            speak("Trying alternative update method")
            try:
                subprocess.run(["UsoClient", "StartScan"], check=True)
                subprocess.run(["UsoClient", "StartDownload"], check=True)
                subprocess.run(["UsoClient", "StartInstall"], check=True)
                speak("Windows update initiated successfully")
                return True
            except:
                speak("Please run Windows Update manually from Settings. Administrator access is required.")
                return False
        
    except Exception as e:
        print(f"Windows update error: {e}")
        speak("Could not complete Windows update. Please try manually from Settings.")
        return False


# ==================== DRIVER UPDATE ====================

def update_drivers():
    """
    Update system drivers using Windows Update
    """
    try:
        if platform.system() != "Windows":
            speak("Driver update is only available on Windows")
            return False
        
        speak("Checking for driver updates")
        print("\n" + "="*60)
        print("DRIVER UPDATE")
        print("="*60)
        
        # Using PowerShell to update drivers
        powershell_cmd = """
        $Session = New-Object -ComObject Microsoft.Update.Session
        $Searcher = $Session.CreateUpdateSearcher()
        $Searcher.ServiceID = '7971f918-a847-4430-9279-4a52d1efe18d'
        $Searcher.SearchScope = 1
        $Searcher.ServerSelection = 3
        $Criteria = "IsInstalled=0 and Type='Driver'"
        $SearchResult = $Searcher.Search($Criteria)
        $Updates = $SearchResult.Updates
        
        if ($Updates.Count -eq 0) {
            Write-Output "No driver updates available"
        } else {
            Write-Output "Found $($Updates.Count) driver updates"
            foreach ($Update in $Updates) {
                Write-Output $Update.Title
            }
        }
        """
        
        try:
            result = subprocess.run(
                ["powershell", "-Command", powershell_cmd],
                capture_output=True,
                text=True,
                timeout=120
            )
            
            output = result.stdout
            
            if "No driver updates available" in output:
                speak("All drivers are up to date")
                print("✓ Drivers are up to date")
                return True
            else:
                print(output)
                speak("Driver updates found. Would you like to install them?")
                confirmation = takecommand()
                
                if "yes" in confirmation.lower():
                    speak("Installing driver updates. Please wait.")
                    
                    install_cmd = """
                    $Session = New-Object -ComObject Microsoft.Update.Session
                    $Searcher = $Session.CreateUpdateSearcher()
                    $Searcher.ServiceID = '7971f918-a847-4430-9279-4a52d1efe18d'
                    $Searcher.SearchScope = 1
                    $Searcher.ServerSelection = 3
                    $Criteria = "IsInstalled=0 and Type='Driver'"
                    $SearchResult = $Searcher.Search($Criteria)
                    
                    $Downloader = $Session.CreateUpdateDownloader()
                    $Downloader.Updates = $SearchResult.Updates
                    $Downloader.Download()
                    
                    $Installer = $Session.CreateUpdateInstaller()
                    $Installer.Updates = $SearchResult.Updates
                    $InstallationResult = $Installer.Install()
                    """
                    
                    subprocess.run(
                        ["powershell", "-Command", install_cmd],
                        capture_output=True,
                        text=True,
                        timeout=600
                    )
                    
                    speak("Driver updates completed. A restart may be required.")
                    return True
                else:
                    speak("Driver update cancelled")
                    return False
                    
        except subprocess.TimeoutExpired:
            speak("Driver update is taking longer. Continuing in background.")
            return True
            
    except Exception as e:
        print(f"Driver update error: {e}")
        speak("Could not update drivers automatically. Please use Device Manager.")
        return False


# ==================== APPLICATION UPDATE ====================

def update_applications():
    """
    Update installed applications using winget (Windows Package Manager)
    """
    try:
        if platform.system() != "Windows":
            speak("Application update using winget is only available on Windows")
            return False
        
        speak("Checking for application updates using Windows Package Manager")
        print("\n" + "="*60)
        print("APPLICATION UPDATE")
        print("="*60)
        
        # Check if winget is installed
        try:
            result = subprocess.run(
                ["winget", "--version"],
                capture_output=True,
                text=True,
                timeout=10
            )
            
            if result.returncode != 0:
                speak("Windows Package Manager not found. Please install it from Microsoft Store.")
                return False
                
        except FileNotFoundError:
            speak("Windows Package Manager not installed. Please install App Installer from Microsoft Store.")
            return False
        
        # List upgradable apps
        speak("Scanning for available updates")
        try:
            result = subprocess.run(
                ["winget", "upgrade"],
                capture_output=True,
                text=True,
                timeout=60
            )
            
            output = result.stdout
            print(output)
            
            if "No installed package found" in output or "No applicable update found" in output:
                speak("All applications are up to date")
                print("✓ All apps are up to date")
                return True
            
            # Count updates
            lines = output.split('\n')
            update_count = len([line for line in lines if line.strip() and not line.startswith('-')])
            
            if update_count > 2:  # Excluding header lines
                speak(f"Found updates for {update_count - 2} applications. Would you like to update all?")
                confirmation = takecommand()
                
                if "yes" in confirmation.lower():
                    speak("Updating all applications. This may take several minutes.")
                    
                    # Update all apps
                    result = subprocess.run(
                        ["winget", "upgrade", "--all", "--silent"],
                        capture_output=True,
                        text=True,
                        timeout=1800  # 30 minutes
                    )
                    
                    if result.returncode == 0:
                        speak("All applications updated successfully")
                        print("✓ Applications updated")
                        return True
                    else:
                        speak("Some applications could not be updated")
                        print(f"⚠️ Update output: {result.stdout}")
                        return False
                else:
                    speak("Application update cancelled")
                    return False
            else:
                speak("All applications are up to date")
                return True
                
        except subprocess.TimeoutExpired:
            speak("Application update is taking longer. Continuing in background.")
            return True
            
    except Exception as e:
        print(f"Application update error: {e}")
        speak("Could not update applications")
        return False


# ==================== MICROSOFT STORE APPS UPDATE ====================

def update_store_apps():
    """
    Update Microsoft Store applications
    """
    try:
        if platform.system() != "Windows":
            speak("Store app update is only available on Windows")
            return False
        
        speak("Updating Microsoft Store applications")
        print("\n" + "="*60)
        print("MICROSOFT STORE UPDATE")
        print("="*60)
        
        # Using PowerShell to update Store apps
        powershell_cmd = """
        Get-CimInstance -Namespace "Root\\cimv2\\mdm\\dmmap" -ClassName "MDM_EnterpriseModernAppManagement_AppManagement01" | 
        Invoke-CimMethod -MethodName UpdateScanMethod
        """
        
        try:
            speak("Initiating Microsoft Store update scan")
            
            result = subprocess.run(
                ["powershell", "-Command", powershell_cmd],
                capture_output=True,
                text=True,
                timeout=120
            )
            
            # Alternative method: Open Store and trigger updates
            subprocess.Popen(["ms-windows-store://downloadsandupdates"])
            
            speak("Microsoft Store opened. It will automatically check and update apps.")
            print("✓ Store update initiated")
            
            return True
            
        except Exception as e:
            print(f"PowerShell error: {e}")
            
            # Fallback: Just open the Store
            speak("Opening Microsoft Store. Please check for updates manually.")
            os.startfile("ms-windows-store:")
            return True
            
    except Exception as e:
        print(f"Store update error: {e}")
        speak("Could not update Store apps automatically")
        return False


# ==================== PYTHON PACKAGES UPDATE ====================

def update_python_packages():
    """
    Update all installed Python packages using pip
    """
    try:
        speak("Updating Python packages. This may take a few minutes.")
        print("\n" + "="*60)
        print("PYTHON PACKAGES UPDATE")
        print("="*60)
        
        # Get list of outdated packages
        try:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "list", "--outdated"],
                capture_output=True,
                text=True,
                timeout=60
            )
            
            output = result.stdout
            print(output)
            
            if not output.strip() or "Package" not in output:
                speak("All Python packages are up to date")
                print("✓ Python packages are up to date")
                return True
            
            # Count outdated packages
            lines = output.strip().split('\n')
            package_count = len(lines) - 2  # Excluding header lines
            
            if package_count > 0:
                speak(f"Found {package_count} outdated packages. Would you like to update them?")
                confirmation = takecommand()
                
                if "yes" in confirmation.lower():
                    speak("Updating Python packages")
                    
                    # Update pip itself first
                    print("Updating pip...")
                    subprocess.run(
                        [sys.executable, "-m", "pip", "install", "--upgrade", "pip"],
                        capture_output=True,
                        timeout=120
                    )
                    
                    # Get package names
                    packages = []
                    for line in lines[2:]:  # Skip header lines
                        if line.strip():
                            package_name = line.split()[0]
                            packages.append(package_name)
                    
                    # Update each package
                    updated = 0
                    failed = 0
                    
                    for package in packages:
                        try:
                            print(f"Updating {package}...")
                            result = subprocess.run(
                                [sys.executable, "-m", "pip", "install", "--upgrade", package],
                                capture_output=True,
                                text=True,
                                timeout=180
                            )
                            
                            if result.returncode == 0:
                                updated += 1
                                print(f"✓ {package} updated")
                            else:
                                failed += 1
                                print(f"✗ {package} failed")
                                
                        except Exception as e:
                            failed += 1
                            print(f"✗ {package} error: {e}")
                    
                    speak(f"Updated {updated} packages successfully. {failed} packages failed.")
                    print(f"\n✓ Updated: {updated}")
                    print(f"✗ Failed: {failed}")
                    return True
                else:
                    speak("Python package update cancelled")
                    return False
            else:
                speak("All Python packages are up to date")
                return True
                
        except subprocess.TimeoutExpired:
            speak("Package update is taking longer. Please check manually.")
            return False
            
    except Exception as e:
        print(f"Python package update error: {e}")
        speak("Could not update Python packages")
        return False


# ==================== COMPLETE SYSTEM UPDATE ====================

def complete_system_update():
    """
    Perform complete system update - all components
    """
    try:
        speak("Starting complete system update. This will take significant time. Please be patient.")
        print("\n" + "="*60)
        print("COMPLETE SYSTEM UPDATE")
        print("="*60 + "\n")
        
        results = {
            "Windows": False,
            "Drivers": False,
            "Applications": False,
            "Store Apps": False,
            "Python Packages": False
        }
        
        # 1. Windows Update
        speak("Step 1 of 5: Updating Windows")
        results["Windows"] = update_windows()
        
        # 2. Driver Update
        speak("Step 2 of 5: Updating Drivers")
        results["Drivers"] = update_drivers()
        
        # 3. Applications Update
        speak("Step 3 of 5: Updating Applications")
        results["Applications"] = update_applications()
        
        # 4. Store Apps Update
        speak("Step 4 of 5: Updating Microsoft Store Apps")
        results["Store Apps"] = update_store_apps()
        
        # 5. Python Packages Update
        speak("Step 5 of 5: Updating Python Packages")
        results["Python Packages"] = update_python_packages()
        
        # Summary
        print("\n" + "="*60)
        print("UPDATE SUMMARY")
        print("="*60)
        
        for component, status in results.items():
            status_icon = "✓" if status else "✗"
            status_text = "Success" if status else "Failed/Skipped"
            print(f"{status_icon} {component}: {status_text}")
        
        successful = sum(results.values())
        total = len(results)
        
        speak(f"System update complete. {successful} out of {total} components updated successfully.")
        
        if successful == total:
            speak("All components updated successfully. A system restart is recommended.")
        else:
            speak("Some components could not be updated. Please check manually.")
        
        return results
        
    except Exception as e:
        print(f"Complete update error: {e}")
        speak("Could not complete full system update")
        return False


# ==================== CHECK FOR UPDATES ====================

def check_for_updates():
    """
    Check for available updates without installing
    """
    try:
        speak("Checking for available updates across all components")
        print("\n" + "="*60)
        print("UPDATE CHECK")
        print("="*60 + "\n")
        
        # Windows Updates
        print("1. Windows Updates:")
        try:
            result = subprocess.run(
                ["powershell", "-Command", "Get-WindowsUpdate"],
                capture_output=True,
                text=True,
                timeout=60
            )
            if "No updates available" in result.stdout:
                print("   ✓ Windows is up to date")
            else:
                print("   ⚠️ Windows updates available")
        except:
            print("   ⚠️ Could not check Windows updates")
        
        # Application Updates
        print("\n2. Application Updates (winget):")
        try:
            result = subprocess.run(
                ["winget", "upgrade"],
                capture_output=True,
                text=True,
                timeout=60
            )
            if "No applicable update found" in result.stdout:
                print("   ✓ All applications up to date")
            else:
                lines = result.stdout.split('\n')
                update_count = len([l for l in lines if l.strip() and not l.startswith('-')])
                print(f"    {update_count - 2} application updates available")
        except:
            print("    Could not check application updates")
        
        # Python Packages
        print("\n3. Python Packages:")
        try:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "list", "--outdated"],
                capture_output=True,
                text=True,
                timeout=60
            )
            lines = result.stdout.strip().split('\n')
            if len(lines) <= 2:
                print("   ✓ All Python packages up to date")
            else:
                print(f"   ⚠️ {len(lines) - 2} Python packages outdated")
        except:
            print("   ⚠️ Could not check Python packages")
        
        speak("Update check complete. Review the results above.")
        return True
        
    except Exception as e:
        print(f"Update check error: {e}")
        speak("Could not complete update check")
        return False


@eel.expose
def allCommands(message=1):
    """Main command processor with auto-continue listening"""
    global CURRENT_LANGUAGE, CURRENT_LANG_CODE, AUTO_LISTEN_MODE, LISTENING_ACTIVE
    
    if message == 1:
        query = takecommand()
        print(query)
        eel.senderText(query)
    else:
        query = message
        eel.senderText(query)
    
    if not query or not query.strip():
        if AUTO_LISTEN_MODE:
            print("No query detected, continuing to listen...")
            allCommands()
        return
    
    try:
        if any(phrase in query.lower() for phrase in ["stop speaking", "stop talking", "be quiet", "shut up", "stop it", "silence"]):
            SPEECH_INTERRUPTED = True
            print("Stop command detected - interrupting speech")
            try:
                pygame.mixer.music.stop()
                pygame.mixer.music.unload()
            except:
                pass
            time.sleep(0.3)
            speak("Stopped")
            if AUTO_LISTEN_MODE:
                allCommands()
            return
        
        elif "stop listening" in query or "stop auto listen" in query or "sleep mode" in query:
            AUTO_LISTEN_MODE = False
            speak("Auto listening disabled. Say IRA or click the mic button to activate me.")
            eel.ShowHood()
            return
        
        elif "start listening" in query or "start auto listen" in query or "wake up" in query:
            AUTO_LISTEN_MODE = True
            speak("Auto listening enabled. I'll keep listening after each command.")
            allCommands()
            return
        # Language commands
        if "change language" in query or "switch language" in query:
            speak("Which language would you like?")
            lang_query = takecommand()
            
            for lang_name in LANGUAGES.keys():
                if lang_name in lang_query:
                    change_language(lang_name)
                    return
            
            speak("Language not recognized. Please try again")
        
        elif "list languages" in query or "available languages" in query:
            list_available_languages()
        
        elif "translate" in query:
            speak("What would you like me to translate?")
            text_to_translate = takecommand()
            
            speak("To which language?")
            target_lang_query = takecommand()
            
            for lang_name in LANGUAGES.keys():
                if lang_name in target_lang_query:
                    target_lang = LANGUAGES[lang_name]["tts"]
                    translated = translate_text(text_to_translate, target_lang)
                    speak(f"Translation: {translated}", language=lang_name)
                    return
        elif "medical" in query or "health" in query or "disease" in query or "बीमारी" in query or "स्वास्थ्य" in query:
            if "about" in query or "information" in query or "ke baare mein" in query or "जानकारी" in query:
                # Extract condition
                condition = query
                for phrase in ["medical information about", "health information about", 
                              "tell me about", "information about", "ke baare mein", "जानकारी"]:
                    condition = condition.replace(phrase, "")
                condition = condition.strip()
                
                # Detect language
                lang = "hindi" if any(hindi_word in query for hindi_word in ["बीमारी", "स्वास्थ्य", "जानकारी"]) else "english"
                
                if condition:
                    get_medical_info(condition, language=lang)
                else:
                    if lang == "hindi":
                        speak("किस बीमारी के बारे में जानकारी चाहिए?", language="hindi")
                    else:
                        speak("What condition would you like information about?")
                    condition = takecommand()
                    if condition:
                        get_medical_info(condition, language=lang)
            else:
                speak("Please specify what medical information you need")
        elif "start continuous listening" in query or "always listen" in query:
           start_continuous_listening()
           speak("Continuous listening activated. I'm always listening now.")
  
        elif "stop continuous listening" in query or "stop always listening" in query:
          stop_continuous_listening()
          speak("Continuous listening stopped. Say IRA to activate me.")
        # WORD DICTATION COMMANDS

        elif"start dictation" in query or "dictate to word" in query or "open word and dictate" in query:
            dictate_to_word_live()

        elif "dictation with punctuation" in query or "punctuation dictation" in query:
            dictate_with_punctuation()

        elif "simple dictation" in query or "dictate now" in query:
            simple_dictation()
        # System commands
        elif "open" in query or "launch" in query or "start" in query:
            from engine.features import openCommand
            openCommand(query)
        
        elif "on youtube" in query:
            from engine.features import PlayYoutube
            PlayYoutube(query)
        elif "change wallpaper" in query or "set wallpaper" in query or "new wallpaper" in query:
            if "from file" in query or "from image" in query:
                speak("Please provide the full path to the image file")
                file_path = takecommand()
                if file_path:
                    change_wallpaper_from_file(file_path)
                else:
                    speak("File path not recognized")
            
            elif "random" in query:
                # Extract category if specified
                categories = ["nature", "space", "abstract", "minimal", "city", "technology"]
                category = "random"
                for cat in categories:
                    if cat in query:
                        category = cat
                        break
                change_wallpaper_random(category)
            
            elif "nature" in query or "forest" in query or "mountain" in query:
                change_wallpaper_random("nature")
            
            elif "space" in query or "galaxy" in query or "stars" in query:
                change_wallpaper_random("space")
            
            elif "abstract" in query or "art" in query:
                change_wallpaper_random("abstract")
            
            elif "minimal" in query or "simple" in query:
                change_wallpaper_random("minimal")
            
            elif "city" in query or "urban" in query:
                change_wallpaper_random("city")
            
            elif "technology" in query or "tech" in query:
                change_wallpaper_random("technology")
            
            elif "generated" in query or "ira" in query:
                change_wallpaper_from_ira_images()
            
            else:
                speak("What kind of wallpaper would you like? Say nature, space, abstract, minimal, city, technology, or random")
                category = takecommand()
                if category:
                    change_wallpaper_random(category)
        
        elif "wallpaper of" in query or "wallpaper with" in query:
            # Search-based wallpaper
            search_query = query.replace("change wallpaper of", "").replace("set wallpaper of", "")
            search_query = search_query.replace("change wallpaper with", "").replace("set wallpaper with", "")
            search_query = search_query.replace("wallpaper of", "").replace("wallpaper with", "").strip()
            
            if search_query:
                change_wallpaper_search(search_query)
            else:
                speak("What should I search for?")
                search_query = takecommand()
                if search_query:
                    change_wallpaper_search(search_query)
        
        elif "restore wallpaper" in query or "default wallpaper" in query:
            restore_default_wallpaper()
        
        elif "random wallpaper" in query:
            change_wallpaper_random("random")
        elif "change wallpaper" in query or "set wallpaper" in query or "new wallpaper" in query:
           if "nature" in query:
             change_wallpaper_random("nature")
           elif "ocean" in query or "sea" in query:
              change_wallpaper_random("ocean")
           elif "space" in query or "galaxy" in query:
              change_wallpaper_random("space")
           elif "sunset" in query:
              change_wallpaper_random("sunset")
           elif "minimal" in query:
              change_wallpaper_random("minimal")
           elif "fire" in query:
               change_wallpaper_random("fire")
           elif "random" in query:
              change_wallpaper_random("random")
           else:
             speak("What kind of wallpaper? Nature, ocean, space, sunset, minimal, or random?")
             category = takecommand()
             if category:
              change_wallpaper_random(category)

        elif "wallpaper of" in query or "wallpaper with" in query:
          search_query = query.replace("wallpaper of", "").replace("wallpaper with", "").strip()
          if search_query:
           change_wallpaper_search(search_query)

        elif "beautiful wallpaper" in query:
           categories = ["nature", "ocean", "sunset", "space"]
           change_wallpaper_random(random.choice(categories))
        elif "system information" in query or "system info" in query:
            get_system_info()
        
        elif "check for updates" in query or "check updates" in query:
            check_for_updates()
        
        elif "update windows" in query or "windows update" in query:
            update_windows()
        
        elif "update drivers" in query or "driver update" in query:
            update_drivers()
        
        elif "update apps" in query or "update applications" in query:
            update_applications()
        
        elif "update store apps" in query or "update microsoft store" in query:
            update_store_apps()
        
        elif "update python packages" in query or "update pip packages" in query:
            update_python_packages()
        
        elif "complete system update" in query or "full system update" in query or "update everything" in query:
            complete_system_update()
        
        elif "update system" in query or "update my system" in query:
            speak("Would you like a complete system update or specific component?")
            choice = takecommand()
            
            if "complete" in choice or "full" in choice or "everything" in choice:
                complete_system_update()
            elif "windows" in choice:
                update_windows()
            elif "driver" in choice:
                update_drivers()
            elif "app" in choice:
                update_applications()
            elif "python" in choice:
                update_python_packages()
            elif "store" in choice:
                update_store_apps()
            else:
                speak("Please specify: windows, drivers, applications, python packages, or complete update")
        elif "weather wallpaper" in query:
            speak("Please tell me the city name")
            city_name = takecommand()
            if city_name:
                get_weather(city_name)
            else:
                speak("City name not recognized")
        elif "change theme" in query or "set theme" in query or "switch theme" in query:
            if "dark" in query or "black" in query or "night" in query:
              change_windows_theme("dark")

            elif "light" in query or "normal" in query or "white" in query or "day" in query:
             change_windows_theme("light")

            else:
             speak("Would you like dark theme or light theme?")
             theme_choice = takecommand()
             if theme_choice:
              if "dark" in theme_choice or "black" in theme_choice:
                change_windows_theme("dark")
              elif "light" in theme_choice or "normal" in theme_choice:
                change_windows_theme("light")
              else:
                speak("Please say dark or light")
        elif "play music" in query or "start music" in query or "resume music" in query:
            play_music()
        
        elif "pause music" in query or "pause the music" in query or "stop the music" in query:
            pause_music()
        
        elif "next song" in query or "next track" in query or "skip song" in query or "skip track" in query:
            next_track()
        
        elif "previous song" in query or "previous track" in query or "last song" in query:
            previous_track()
        
        elif "stop music" in query or "stop playing" in query:
            stop_music()
        
        # OPEN MUSIC APPS
        
        elif "open spotify" in query or "launch spotify" in query or "start spotify" in query:
            open_spotify()
        
        elif "play spotify" in query or "spotify play" in query:
            if "play" in query and any(word in query for word in ["song", "track", "music"]):
                # Extract song name
                song_name = query.replace("play", "").replace("on spotify", "").replace("spotify", "")
                song_name = song_name.replace("song", "").replace("track", "").strip()
                if song_name:
                    play_on_spotify(song_name)
                else:
                    open_spotify()
            else:
                open_spotify()
        
        elif "youtube music" in query:
            if "play" in query:
                song_name = query.replace("play", "").replace("on youtube music", "").replace("youtube music", "")
                song_name = song_name.strip()
                if song_name:
                    play_on_youtube_music(song_name)
                else:
                    speak("What would you like to play?")
                    song_name = takecommand()
                    if song_name:
                        play_on_youtube_music(song_name)
        
        # PLAY SPECIFIC SONGS
        
        elif "play song" in query or "play track" in query:
            # Extract song name
            song_name = query.replace("play song", "").replace("play track", "").replace("play", "")
            song_name = song_name.replace("on spotify", "").replace("on youtube", "").strip()
            
            if song_name:
                # Ask which platform
                speak(f"Where would you like to play {song_name}? Spotify, YouTube Music, or local?")
                platform = takecommand()
                
                if "spotify" in platform:
                    play_on_spotify(song_name)
                elif "youtube" in platform:
                    play_on_youtube_music(song_name)
                elif "local" in platform:
                    play_local_music(song_name)
                else:
                    # Default to YouTube Music (free)
                    play_on_youtube_music(song_name)
            else:
                speak("What song would you like to play?")
                song_name = takecommand()
                if song_name:
                    play_on_youtube_music(song_name)
        
        # LOCAL MUSIC PLAYBACK
        
        elif "play local music" in query or "play from library" in query:
            if "song" in query or "track" in query:
                song_name = query.replace("play local music", "").replace("play from library", "")
                song_name = song_name.replace("song", "").replace("track", "").strip()
                play_local_music(song_name)
            else:
                play_local_music()
        
        elif "scan music library" in query or "scan music" in query or "refresh music library" in query:
            scan_local_music_library()
        
        # PLAYLIST MANAGEMENT
        
        elif "create playlist" in query or "make playlist" in query or "new playlist" in query:
            speak("What should I name the playlist?")
            playlist_name = takecommand()
            if playlist_name:
                create_playlist(playlist_name)
            else:
                speak("Playlist name not recognized")
        
        elif "add to playlist" in query or "add song to playlist" in query:
            speak("Which playlist?")
            playlist_name = takecommand()
            
            if playlist_name:
                speak("Which song would you like to add?")
                song_name = takecommand()
                
                if song_name:
                    add_to_playlist(playlist_name, song_name)
                else:
                    speak("Song name not recognized")
            else:
                speak("Playlist name not recognized")
        
        elif "play playlist" in query:
            playlist_name = query.replace("play playlist", "").strip()
            
            if not playlist_name:
                speak("Which playlist would you like to play?")
                playlist_name = takecommand()
            
            if playlist_name:
                play_playlist(playlist_name)
            else:
                speak("Playlist name not recognized")
        
        elif "list playlists" in query or "show playlists" in query or "my playlists" in query:
            list_playlists()
        
        # SHUFFLE AND REPEAT
        
        elif "shuffle" in query:
            if "on" in query or "enable" in query or "turn on" in query:
                shuffle_mode(True)
            elif "off" in query or "disable" in query or "turn off" in query:
                shuffle_mode(False)
            else:
                shuffle_mode(True)
        
        elif "repeat" in query:
            if "one" in query or "this" in query or "song" in query:
                repeat_mode("one")
            elif "all" in query or "playlist" in query:
                repeat_mode("all")
            elif "off" in query or "disable" in query:
                repeat_mode("off")
            else:
                repeat_mode("all")
        
        # MUSIC INFO
        
        elif "what song is this" in query or "what's playing" in query or "whats playing" in query or "current song" in query or "now playing" in query:
            what_is_playing()
        
        elif "what music is playing" in query or "what's this song" in query:
            what_is_playing()
        
        # MUSIC VOLUME CONTROL
        
        elif "music volume" in query or "set music volume" in query:
            # Extract volume level
            import re
            numbers = re.findall(r'\d+', query)
            
            if numbers:
                volume_level = int(numbers[0])
                music_volume_control("set", volume_level)
            elif "up" in query or "increase" in query or "louder" in query:
                music_volume_control("increase")
            elif "down" in query or "decrease" in query or "lower" in query:
                music_volume_control("decrease")
            else:
                speak("What volume level? Say a number from 0 to 100")
                volume_text = takecommand()
                try:
                    volume_level = int(''.join(filter(str.isdigit, volume_text)))
                    music_volume_control("set", volume_level)
                except:
                    speak("Invalid volume level")
        
        elif "louder" in query or "increase music volume" in query:
            music_volume_control("increase")
        
        elif "quieter" in query or "decrease music volume" in query or "lower music volume" in query:
            music_volume_control("decrease")
        
        # QUICK MUSIC COMMANDS
        
        elif query in ["play", "pause", "next", "previous", "skip"]:
            if query == "play":
                play_music()
            elif query == "pause":
                pause_music()
            elif query in ["next", "skip"]:
                next_track()
            elif query == "previous":
                previous_track()
        elif "dark mode" in query or "enable dark mode" in query or "turn on dark mode" in query:
          change_windows_theme("dark")

        elif "light mode" in query or "enable light mode" in query or "turn on light mode" in query:
          change_windows_theme("light")

        elif "toggle theme" in query or "switch theme mode" in query:
          toggle_theme()

        elif "current theme" in query or "what theme" in query or "which theme" in query:
           get_current_theme()

# ACCENT COLOR COMMANDS
        elif "change accent color" in query or "set accent color" in query or "accent color" in query:
           colors = ["blue", "red", "green", "purple", "orange", "pink", "gray", "yellow", "teal"]

           color_found = None
           for color in colors:
             if color in query:
              color_found = color
              break

             if color_found:
              change_accent_color(color_found)
             else:
              speak("What color would you like? Available colors are: blue, red, green, purple, orange, pink, gray, yellow, teal")
              color_choice = takecommand()
              if color_choice:
                for color in colors:
                  if color in color_choice:
                    change_accent_color(color)
                    break

# TRANSPARENCY COMMANDS
        elif "enable transparency" in query or "turn on transparency" in query:
          enable_transparency(True)

        elif "disable transparency" in query or "turn off transparency" in query:
          enable_transparency(False)

# THEME PRESET COMMANDS
        elif "apply preset" in query or "theme preset" in query or "apply theme preset" in query:
          presets = ["blue_dark", "nature_light", "sunset", "ocean", "minimal", "cyberpunk"]

          preset_found = None
          for preset in presets:
            preset_words = preset.replace("_", " ")
            if preset_words in query or preset in query:
              preset_found = preset
              break

          if preset_found:
            apply_custom_theme_preset(preset_found)
          else:
            speak("Which preset would you like? Say: blue dark, nature light, sunset, ocean, minimal, or cyberpunk")
            preset_choice = takecommand()
            if preset_choice:
              preset_choice = preset_choice.replace(" ", "_")
              apply_custom_theme_preset(preset_choice)

        elif "list presets" in query or "available presets" in query or "show presets" in query:
          list_available_presets()

# QUICK THEME PRESETS
        elif "blue dark theme" in query or "blue dark preset" in query:
          apply_custom_theme_preset("blue_dark")

        elif "nature light theme" in query or "nature light preset" in query:
          apply_custom_theme_preset("nature_light")

        elif "sunset theme" in query or "sunset preset" in query:
            apply_custom_theme_preset("sunset")

        elif "ocean theme" in query or "ocean preset" in query:
            apply_custom_theme_preset("ocean")

        elif "minimal theme" in query or "minimal preset" in query:
           apply_custom_theme_preset("minimal")

        elif "cyberpunk theme" in query or "cyberpunk preset" in query:
           apply_custom_theme_preset("cyberpunk")

        elif "shutdown" in query or "turn off" in query:
            shutdown_pc()
        elif "restart" in query or "reboot" in query:
            restartthepc()
        elif "lock" in query:
            lock_pc()
        
        elif "screenshot" in query or "capture screen" in query:
            take_screenshot()
        
        elif "increase volume" in query:
            increase_volume()
        
        elif "decrease volume" in query:
            decrease_volume()
        
        elif "mute" in query:
            mute_volume()
        
        elif "internet speed" in query or "speed test" in query:
            check_internet_speed()
        
        elif "time" in query and "current" in query:
            tell_time()
        
        elif "date" in query and ("current" in query or "today" in query):
            tell_date()
        
        elif "create note" in query or "take note" in query:
            speak("What would you like me to write in the note?")
            note_text = takecommand()
            if note_text:
                create_note(note_text)
            else:
                speak("Note text not recognized")
        
        elif "create folder" in query or "make folder" in query:
            speak("What should be the folder name?")
            folder_name = takecommand()
            if folder_name:
                create_folder(folder_name)
            else:
                speak("Folder name not recognized")
        
        elif "delete file" in query or "remove file" in query:
            speak("Please tell me the full path of the file to delete")
            file_path = takecommand()
            file_path = file_path.replace(" slash ", "/").replace(" backslash ", "\\")
            if file_path:
                delete_file(file_path)
            else:
                speak("File path not recognized")
        elif "increase brightness" in query or "decrease brightness" in query or "set brightness" in query:
            adjust_brightness(query)
        elif "battery" in query:
            check_battery_status()
        
        elif "sleep" in query and "pc" in query:
            sleepthepc()
        
        elif "empty recycle bin" in query or "clear recycle bin" in query:
            empty_recycle_bin()
        
        elif "wikipedia" in query or "search wikipedia" in query:
            speak("What topic should I search for?")
            topic = takecommand()
            if topic:
                search_wikipedia(topic)
            else:
                speak("Topic not recognized.")
        
        elif "add task" in query or "add to-do" in query:
            speak("What task would you like to add?")
            task = takecommand()
            if task:
                add_task(task)
            else:
                speak("Task not recognized.")
        
        elif "show tasks" in query or "read tasks" in query or "my tasks" in query:
            show_tasks()
        
        elif "delete task" in query or "remove task" in query:
            speak("Please tell me the task number to remove")
            task_number_text = takecommand()
            try:
                task_number = int(task_number_text)
                delete_task(task_number)
            except ValueError:
                speak("Invalid task number. Please say a number.")
        
        elif "set reminder" in query or "remind me" in query:
            speak("What should I remind you about?")
            task = takecommand()
            speak("At what time? Please say in HH:MM 24-hour format")
            time_str = takecommand()
            time_str = time_str.replace(" ", "").replace(".", ":")
            set_reminder(task, time_str)
        
        elif "disk usage" in query or "disk space" in query or "free space" in query:
            speak("Which drive would you like me to check?")
            drive_letter = takecommand().strip().lower()
            drive_map = {"see": "C", "sea": "C", "c": "C", "d": "D", "dee": "D"}
            drive = drive_map.get(drive_letter, drive_letter.upper())
            check_disk_usage(drive)
        elif "check cpu" in query or "cpu usage" in query or "ram usage" in query or "system usage" in query:
            check_system_usage()
        elif "wifi info" in query or "wi-fi info" in query or "network name" in query:
            get_wifi_info()
        
        elif "network info" in query or "ip address" in query or "my ip" in query:
            get_network_info()
        
        elif "wifi password" in query or "wi-fi password" in query or "network password" in query:
            if "current" in query or "this" in query:
                show_wifi_password()
            else:
                speak("Which network password would you like?")
                network_name = takecommand()
                if network_name:
                    show_wifi_password(network_name)
        
        elif "scan wifi" in query or "available networks" in query or "list networks" in query:
            list_available_wifi()
        elif "who are you" in query or "what is your name" in query or "your name" in query:
            speak("I am IRA, your intelligent voice assistant. I'm here to help you with various tasks.")
        
        elif "who created you" in query or "who made you" in query or "your creator" in query or "who developed you" in query:
            speak("I was created by Ankita, Anjali, Shubham and Amrita.")
        
        elif "who is your creator" in query or "who is your developer" in query:
            speak("My creators are Ankita, Anjali, Shubham and Amrita. They developed me to assist you.")
        
        elif "tell me about yourself" in query or "introduce yourself" in query:
            speak("Hello! I am IRA, an intelligent voice assistant created by Ankita, Anjali, Shubham and Amrita. I can help you with many tasks like checking weather, managing files, controlling your system, and much more. Just ask me anything!")

        elif "what is unique" in query or "what makes you different" in query or "what makes you special" in query or "your uniqueness" in query or "why choose you" in query or "your specialty" in query:
                        unique_features = """What makes me truly unique? I speak 27 plus languages fluently, including all major Indian languages. I can read files just by their name, no path needed. I generate images, create presentations, and write complete code in any language. I reveal WiFi passwords, control your entire system, and provide news in both English and Hindi. Most importantly, I keep listening automatically, so you never have to activate me repeatedly.I'm always ready to help in your preferred language."""
                        speak(unique_features)
        elif "read file" in query or "read text" in query or "फ़ाइल पढ़ो" in query:
            speak("What is the file name?")
            filename = takecommand()
    
            if filename:
               speak("Should I read in English or Hindi?")
               lang_choice = takecommand()
        
               lang = "hindi" if "hindi" in lang_choice or "हिंदी" in lang_choice else "english"
               read_text_file_smart(filename, language=lang)
            else:
               speak("File name not recognized")

        elif "read pdf" in query or "पीडीएफ पढ़ो" in query:
            speak("What is the PDF file name?")
            filename = takecommand()
    
            if filename:
              speak("Should I read in English or Hindi?")
              lang_choice = takecommand()
        
              lang = "hindi" if "hindi" in lang_choice or "हिंदी" in lang_choice else "english"
        
              speak("Should I read all pages or specific pages?")
              page_choice = takecommand()
        
              if "specific" in page_choice:
                speak("From which page?")
                start_text = takecommand()
                try:
                  start_page = int(''.join(filter(str.isdigit, start_text))) - 1
                except:
                  start_page = 0
            
                speak("To which page?")
                end_text = takecommand()
                try:
                  end_page = int(''.join(filter(str.isdigit, end_text)))
                except:
                  end_page = None
            
                read_pdf_file_smart(filename, start_page, end_page, language=lang)
              else:
                 read_pdf_file_smart(filename, language=lang)
            else:
              speak("File name not recognized")

        elif "read word" in query or "read document" in query or "वर्ड पढ़ो" in query:
             speak("What is the Word document name?")
             filename = takecommand()
    
             if filename:
               speak("Should I read in English or Hindi?")
               lang_choice = takecommand()
        
               lang = "hindi" if "hindi" in lang_choice or "हिंदी" in lang_choice else "english"
               read_word_document_smart(filename, language=lang)
             else:
               speak("File name not recognized")

        elif "hindi news" in query or "हिंदी समाचार" in query or "hindi mein news" in query:
          get_news_hindi_from_rss(language="hindi")
    
        elif "news in hindi" in query or "समाचार सुनाओ" in query:
           get_news_hindi_from_rss(language="hindi")
        elif "news" in query or "current affairs" in query or "headlines" in query:
           if "business" in query:
             get_business_news()
           elif "sports" in query:
             get_sports_news()
           elif "technology" in query or "tech" in query:
             get_tech_news()
           elif "entertainment" in query:
             get_entertainment_news()
           else:
             get_current_affairs()

        elif "search news" in query or "news about" in query:
           topic = query.replace("search news", "").replace("news about", "")
           topic = topic.replace("for", "").strip()
    
           if topic:
            search_news(topic)
           else:
            speak("What topic should I search for?")
            topic = takecommand()
            if topic:
              search_news(topic)  
        elif "send message" in query or "phone call" in query or "video call" in query:
           from engine.features import findContact, whatsApp, makeCall, sendMessage
           contact_no, name = findContact(query)
    
           if contact_no != 0:
            speak("Which mode you want to use whatsapp or mobile")
            preference = takecommand()

           if "mobile" in preference:
              if "send message" in query or "send sms" in query: 
                 speak("What message to send")
                 message = takecommand()
                 sendMessage(message, contact_no, name)
              elif "phone call" in query:
                makeCall(name, contact_no)
              else:
                speak("Please try again")
        
           elif "whatsapp" in preference:
            if "send message" in query:
                speak("What message to send")
                message_text = takecommand()  # Changed variable name
                whatsApp(contact_no, message_text, 'message', name)  # Fixed parameters
            elif "phone call" in query:
                whatsApp(contact_no, '', 'call', name)  # Fixed parameters
            else:  # video call
                whatsApp(contact_no, '', 'video call', name)  # Fixed parameters
           else:
            speak("I didn't understand. Please say whatsapp or mobile")
        elif "generate image" in query or "create image" in query or "draw image" in query or "make image" in query:
            speak("Please describe the image you want in detail....")
            image_prompt = takecommand()

            if image_prompt and len(image_prompt) > 3:
            # Use auto mode for best results
               generate_image(image_prompt, service="auto")
            else:
               speak("I didn't catch the description. Please try again.")
            
        elif "generate photo" in query or "create picture" in query or "make picture" in query:
            speak("Describe the photo you want me to create..")
            image_prompt = takecommand()

            if image_prompt and len(image_prompt) > 3:
                generate_image(image_prompt, service="auto")
            else:
                speak("Image description not clear.")

        # PowerPoint Generation Commands
        elif "create presentation" in query or "make presentation" in query or "create ppt" in query or "make ppt" in query:
            speak("Would you like to create from a template or from a topic?")
            choice = takecommand()
            
            if "template" in choice:
                speak("Please provide the full path to your template file")
                template_path = takecommand()
                speak("What should be the presentation name?")
                ppt_name = takecommand()
                
                if template_path and ppt_name:
                    create_ppt_from_template(template_path, ppt_name)
                else:
                    speak("Template path or name not recognized")
            
            elif "topic" in choice:
                speak("What topic should the presentation be about?")
                topic = takecommand()
                
                if topic:
                    speak("How many slides would you like? Say a number between 5 and 15")
                    num_slides_text = takecommand()
                    
                    try:
                        num_slides = int(''.join(filter(str.isdigit, num_slides_text)))
                        num_slides = max(5, min(num_slides, 15))
                    except:
                        num_slides = 7
                        speak("Using default 7 slides")
                    
                    speak("Which design theme? Professional, modern, or dark?")
                    theme_choice = takecommand()
                    theme = "professional"
                    if "modern" in theme_choice:
                        theme = "modern"
                    elif "dark" in theme_choice:
                        theme = "dark"
                    
                    create_ppt_from_topic(topic, num_slides, theme)
                else:
                    speak("Topic not recognized")
            else:
                speak("Please specify template or topic")
        
        elif "presentation on" in query or "ppt on" in query:
            # Quick presentation creation
            topic = query.replace("create presentation on", "").replace("make presentation on", "")
            topic = topic.replace("create ppt on", "").replace("make ppt on", "").strip()
            
            if topic:
                create_ppt_from_topic(topic, num_slides=7, design_theme="professional")
            else:
                speak("Please specify the topic")
        #  CODE GENERATION COMMANDS 
        elif "write code" in query or "generate code" in query or "create code" in query:
            speak("What kind of code should I write?")
            code_description = takecommand()
            
            if code_description and len(code_description) > 3:
                open_vscode_and_write_code(code_description)
            else:
                speak("I didn't catch that. Please describe the code you need.")

        elif "code for" in query or "program for" in query or "script for" in query:
            # Extract description from query
            description = query
            for phrase in ["write code for", "generate code for", "create code for", 
                          "code for", "program for", "script for", "write program for"]:
                description = description.replace(phrase, "")
            description = description.strip()
            
            if description:
                open_vscode_and_write_code(description)
            else:
                speak("Please describe what the code should do")

        elif "open vs code" in query or "open vscode" in query or "launch vs code" in query:
            speak("Opening Visual Studio Code")
            try:
                subprocess.Popen(['code'])
                speak("VS Code opened")
            except:
                speak("Could not open VS Code. Please check if it's installed.")
        else:
            if query.strip():
                from engine.features import geminiai
                geminiai(query)

        
    except Exception as e:
        print(f"Error: {e}")
        speak("Sorry, I encountered an error")
    
    if AUTO_LISTEN_MODE:
        print("\n🔄 Ready for next command...")
        time.sleep(0.3)  # Brief pause
        allCommands()  # This keeps siriwave active and continues listening
    else:
        eel.ShowHood()

@eel.expose
def stop_speech():
    """Stop current speech - callable from frontend button"""
    global SPEECH_INTERRUPTED
    SPEECH_INTERRUPTED = True
    print("Speech stopped by frontend button")
    try:
        pygame.mixer.music.stop()
        pygame.mixer.music.unload()
    except:
        pass
    speak("Stopped")
def hotword_listener():
    """Listen for wake word 'IRA' and auto-start listening"""
    global AUTO_LISTEN_MODE
    
    r = sr.Recognizer()
    r.energy_threshold = 3000
    r.dynamic_energy_threshold = True
    r.pause_threshold = 0.6
    
    with sr.Microphone() as source:
        print("👂 Listening for 'IRA'...")
        r.adjust_for_ambient_noise(source, duration=0.3)
        
        try:
            audio = r.listen(source, timeout=None, phrase_time_limit=2)
            query = r.recognize_google(audio, language='en-in').lower()
            
            if "ira" in query:
                print("✅ Wake word detected!")
                AUTO_LISTEN_MODE = True
                speak("Yes, how can I help?")
                allCommands()
                
        except sr.UnknownValueError:
            pass
        except sr.RequestError as e:
            print(f"Service error: {e}")
            time.sleep(2)
        except Exception as e:
            print(f"Hotword error: {e}")
            time.sleep(0.5)

def start_hotword_thread():
    """Run hotword detection in background"""
    while True:
        try:
            hotword_listener()
        except Exception as e:
            print(f"Hotword thread error: {e}")
            time.sleep(1)

def initialize():
    """Initialize the assistant with auto-listen mode"""
    global AUTO_LISTEN_MODE
    
    print("Initializing Multi-Language Voice Assistant...")
    print("\n=== AUTO-LISTEN MODE ===")
    print("✅ Enabled - I'll keep listening after each command")
    print("Say 'stop listening' to disable")
    print("Say 'start listening' to enable")
    print("Press SPACE key to interrupt speech")
    print("========================\n")
    
    speak("Voice assistant initialized with auto listening enabled")
    print(f"Current language: {CURRENT_LANGUAGE}")
    
    reminder_thread = threading.Thread(target=reminder_checker, daemon=True)
    reminder_thread.start()
    
    hotword_thread = threading.Thread(target=start_hotword_thread, daemon=True)
    hotword_thread.start()
    
    print("\n✅ Ready! Say 'IRA' to start or click the mic button.")


if __name__ == "__main__":
    initialize()
    
    # Start Eel web interface
    eel.init("web")  # Change "web" to your actual web folder name
    eel.start("index.html", mode='chrome', size=(1000, 600), port=8080)